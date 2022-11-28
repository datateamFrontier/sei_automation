import pandas as pd
import numpy as np

import olah_data as od

import pptx
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_TICK_LABEL_POSITION, XL_LABEL_POSITION, XL_TICK_MARK,XL_MARKER_STYLE, XL_LEGEND_POSITION
from pptx.oxml.xmlchemy import OxmlElement


"""
Program ini berisi fungsi-fungsi untuk melakukan plot data hasil olah indeks ke file laporan .pptx (v2)
"""


"""
nama_bulan dan nomor_bulan digunakan sebagai acuan ketika konversi bulan diperlukan
"""
nama_bulan = ['Januari','Februari','Maret','April','Mei','Juni',
              'Juli','Agustus','September','Oktober','November','Desember']
nomor_bulan = ['1','2','3','4','5','6',
              '7','8','9','10','11','12']


def Reverse(lst):
    """
    Membalik urutan list array

    Parameter
    ---------
    lst : list array yang akan dibalik

    Return
    ------
    list array yang sudah terbalik urutannya
    """

    return [ele for ele in reversed(lst)]


def bobot_indeks(bobot_indeks, perbaikan=False):
    """
    Melakukan pra-pengolahan data bobot. Mengisi kolom-kolom yang kosong di bagian 'DIMENSI',
    'KPI', 'SUB KPI', dan 'ASPEK'.

    Parameter
    ---------
    bobot_indeks : dataframe yang berisi daftar bobot

    Return
    ------
    bobot_indeks : dataframe bobot yang sudah tidak memiliki cell kosong di kolom yang ditentukan.
    """
    if perbaikan==False:
        bobot_indeks['DIMENSI'] = bobot_indeks['DIMENSI'].fillna(method='ffill')
        bobot_indeks['KPI'] = bobot_indeks['KPI'].fillna(method='ffill')
    bobot_indeks['SUB KPI'] = bobot_indeks['SUB KPI'].fillna(method='ffill')
    bobot_indeks['ASPEK'] = bobot_indeks['ASPEK'].fillna(method='ffill')

    return bobot_indeks



def tabel_definisi_dimensi(kriteria):
    """
    Memilih tabel definisi sesuai kriteria

    Parameter
    ---------
    kriteria : kriteria SEI ('call center', 'online chat', 'twitter', atau 'email')

    Return
    ------
    isi_tabel : list array yang dipakai untuk tabel definisi dimensi
    """
    if kriteria=='call center':
        isi_tabel_akses = [['DIMENSI : ACCESS',
                      'Mengukur kinerja call center dari segi kemudahan yang dialami pelanggan saat menghubungi call center berdasarkan 3 kpi yaitu : ACCESSIBILITY, AVAILIBILITY, DAN CONNECTION SPEED.'],
                     ['KPI : ACCESSIBILITY',
                      'Mengukur tingkat kecepatan yang dialami pelanggan untuk mendapatkan nada sambung pada saat menghubungi call center.'],
                     ['KPI : AVAILIBILITY',
                      'Mengukur tingkat keberhasilan yang dialami pelanggan untuk terhubung dengan mesin penjawab atau mesin Interactive Voice Response (IVR) atau Call Center Officer (CCO) setelah menekan nomor call center.'],
                     ['KPI : CONNECTION SPEED',
                      'Mengukur tingkat kecepatan yang dialami pelanggan untuk menerima salam pembuka (mulai dilayani) dari mesin penjawab atau IVR setelah terhubung dengan call center atau dari CCO setelah menekan menu "berbicara dengan CCO".']]
        isi_tabel_sispro = [['DIMENSI : SYSTEM & PROCEDURE',
                      'Mengukur kinerja call center dari segi kemudahan dan kenyamanan yang dialami pelanggan selama menggunakan call center berdasarkan 3 KPI yaitu : SERVICE STANDARD, ENJOYING, dan SYSTEM.'],
                     ['KPI : SERVICE STANDARD CONSISTENCY',
                      'Mengukur tingkat profesionalisme pelayanan yang dialami pelanggan sehubungan dengan konsistensi penerapan standar pelayanan saat pelanggan berinteraksi dengan mesin IVR atau CCO.'],
                     ['KPI : ENJOYING',
                      'Mengukur tingkat kenyamanan yang dialami pelanggan saat menggunakan mesin IVR atau berbicara dengan CCO.'],
                     ['KPI : SYSTEM',
                      'Mengukur tingkat kemudahan yang dialami pelanggan untuk mengikuti aliran proses pelayanan dari mesin IVR.']]
        isi_tabel_people = [['DIMENSI : PEOPLE',
                      'Mengukur kinerja call center dari segi keramahtamahan dan kualitas solusi yang dialami pelanggan saat dilayani oleh CCO berdasarkan 2 KPI yaitu : SOFT SKILL dan HARD SKILL.'],
                     ['KPI : SOFT SKILL',
                      'Mengukur tingkat keramahtamahan yang dialami pelanggan saat dilayani CCO.'],
                     ['KPI : HARD SKILL',
                      'Mengukur tingkat kualitas solusi yang didapatkan pelangan atas kebutuhan informasi produk/jasa yang diberikan oleh CCO.']]
        isi_tabel = [isi_tabel_akses, isi_tabel_sispro, isi_tabel_people]

    elif kriteria=='twitter':
        isi_tabel_engaging = [['DIMENSI : ENGAGING',
                      'Mengukur kinerja sebuah akun Twitter dari segi keterikatan yang dialami pelanggan saat bergabung ke akun tersebut dalam hal kemudahan akses dan kemudahan sistem didalamnya berdasarkan 2 KPI yaitu : ACCESS dan SYSTEM.'],
                     ['KPI : ACCESS',
                      'Mengukur tingkat kemudahan yang dialami pelanggan untuk bergabung dalam akun Twitter tersebut didalamnya mengandung unsur mudah ditemukan akun Twitternya, mudah terkirim pesannya dan ada penanggung jawab diakun tersebut.'],
                     ['KPI : SYSTEM',
                      'Mengukur tingkat kepastian pesan yang dikirim pelanggan dapat terjawab/ditanggapi secara langsung, cepat responnya dan tuntas jawabnya.']]
        isi_tabel_humtouch = [['DIMENSI : HUMAN TOUCHING',
                      'Mengukur kinerja sebuah akun Twitter dari segi sentuhan emosional kemanusiaan sehingga pelanggan merasa sedang berhubungan dengan petugas secara langsung berdasarkan 2 KPI yaitu : COMMUNICATING dan ATTITUDE.'],
                     ['KPI : COMMUNICATING',
                      'Mengukur keramahan dan keluwesan dari bahasa yang digunakan, yang dirasakan pelanggan saat berinteraksi baik verbal, non verbal maupun virtual.'],
                     ['KPI : ATTITUDE',
                      'Mengukur kinerja sebuah akun Twitter dari segi perhatian, emphaty dan kepedulian yang dirasakan pelanggan dalam menanggapi masalahnya.']]
        isi_tabel_navigat = [['DIMENSI : NAVIGATING',
                      'Mengukur sebuah akun Twitter dalam hal ketuntasan dan solusi tindak lanjut dari hal yang disampaikan pelanggan berdasarkan 3 KPI yaitu : PROBING, PROVIDING SOLUTION dan CLOSING.'],
                     ['KPI : PROBING',
                      'Mengukur tingkat kedalaman dalam mencari kebutuhan pelanggan yang diharapkan.'],
                     ['KPI : PROVIDING SOLUTION',
                      'Mengukur tingkat ketuntasan dalam memberikan solusi dan menawarkan solusi lanjutan yang dirasakan oleh pelanggan.'],
                     ['KPI : CLOSING',
                      'Mengukur tingkat keramahan dan keterkaitan lebih lanjut untuk menggunakan/berkunjung ke akun Twitter tersebut.']]
        isi_tabel = [isi_tabel_engaging, isi_tabel_humtouch, isi_tabel_navigat]
        
    elif kriteria=='online chat':
        isi_tabel_engaging = [['DIMENSI : ENGAGING',
                      'Mengukur kinerja sebuah akun Online Chat dari segi keterikatan yang dialami pelanggan saat bergabung ke akun tersebut dalam hal kemudahan akses dan kemudahan sistem didalamnya berdasarkan 2 KPI yaitu : ACCESS dan SYSTEM.'],
                     ['KPI : ACCESS',
                      'Mengukur tingkat kemudahan yang dialami pelanggan untuk bergabung dalam akun Online Chat tersebut didalamnya mengandung unsur mudah ditemukan akun Online Chat-nya, mudah terkirim pesannya dan ada penanggung jawab diakun tersebut.'],
                     ['KPI : SYSTEM',
                      'Mengukur tingkat kepastian pesan yang dikirim pelanggan dapat terjawab/ditanggapi secara langsung, cepat responnya dan tuntas jawabnya.']]
        isi_tabel_humtouch = [['DIMENSI : HUMAN TOUCHING',
                      'Mengukur kinerja sebuah akun Online Chat dari segi sentuhan emosional kemanusiaan sehingga pelanggan merasa sedang berhubungan dengan petugas secara langsung berdasarkan 2 KPI yaitu : COMMUNICATING dan ATTITUDE.'],
                     ['KPI : COMMUNICATING',
                      'Mengukur keramahan dan keluwesan dari bahasa yang digunakan, yang dirasakan pelanggan saat berinteraksi baik verbal, non verbal maupun virtual.'],
                     ['KPI : ATTITUDE',
                      'Mengukur kinerja sebuah akun Online Chat dari segi perhatian, emphaty dan kepedulian yang dirasakan pelanggan dalam menanggapi masalahnya.']]
        isi_tabel_navigat = [['DIMENSI : NAVIGATING',
                      'Mengukur sebuah akun Online Chat dalam hal ketuntasan dan solusi tindak lanjut dari hal yang disampaikan pelanggan berdasarkan 3 KPI yaitu : PROBING, PROVIDING SOLUTION dan CLOSING.'],
                     ['KPI : PROBING',
                      'Mengukur tingkat kedalaman dalam mencari kebutuhan pelanggan yang diharapkan.'],
                     ['KPI : PROVIDING SOLUTION',
                      'Mengukur tingkat ketuntasan dalam memberikan solusi dan menawarkan solusi lanjutan yang dirasakan oleh pelanggan.'],
                     ['KPI : CLOSING',
                      'Mengukur tingkat keramahan dan keterkaitan lebih lanjut untuk menggunakan/berkunjung ke akun Online Chat tersebut.']]
        isi_tabel = [isi_tabel_engaging, isi_tabel_humtouch, isi_tabel_navigat]        

    elif kriteria=='email':
        isi_tabel_enab = [['DIMENSI : ENABLING',
                      'Mengukur kinerja tersedianya layanan Email Customer Service bagi pelanggan pengguna Email. Perusahaan Customer Centric membuat pelanggan mampu dan mau menggunakan layanan Email secara tersistem.'],
                     ['KPI : FEASIBILITY',
                      'Mengukur tingkat komunikasi suatu alamat Email Customer Service Serta kemudahan untuk menghafal dan menggunakan alamat Email tersebut.'],
                     ['KPI : ACCESSIBILITY',
                      'Mengukur tingkat kemudahan untuk mengakses alamat Email, tidak adanya kesulitan dalam mengirimkan email ke Customer Service untuk minta dilayani ditandai oleh adanya tindak lanjut.'],
                     ['KPI : AVAIBILITY',
                      'Mengukur tingkat Consistency ketersediaan layanan Email Customer Service termasuk Realibility layanan Email Customer Service yang disediakan.']]
        isi_tabel_enjoy1 = [['DIMENSI : ENJOYING',
                      'Mengukur kinerja total dari beberapa kumpulan indikator yang mengukur kenyamanan pelanggan dalam menggunakan layanan Email Customer Service perusahaan'],
                     ['KPI : SYSTEM & PROCEDURE',
                      '1. Indikator yang mengukur adanya system, konsistensi system serta kemudahan dan kekonsistenan dalam menyajikan kemudahan bagi pelanggan pengguna Email Customer Service.\n'+
                      '2. Evaluasi juga dilakukan terhadap parameter ada tidaknya, efektif tidaknya sistem record dan retrieving data untuk layanan bertingkat.\n'+
                      '3. Keberadaan sistem menindak lanjuti (follow up) permintaan pelanggan dalam Email Customer Service menjadi permasalahan.\n'+
                      '4. Kunci ketuntasan dalam pelayanan Email Customer Service yang juga merupakan indikator dipantau dalam riset ini.\n'+
                      '5. Dalam System and Procedure juga diukur service level yakni kecepatan dan keakurasian dan konsistensi pelayanan Email Customer Service lewat pintu Email.']]
        isi_tabel_enjoy2 = [['DIMENSI : ENJOYING',
                      'Mengukur kinerja total dari beberapa kumpulan indikator yang mengukur kenyamanan pelanggan dalam menggunakan layanan Email Customer Service perusahaan.'],
                     ['KPI : PEOPLE',
                      '1. Berbicara soal people maka Soft Skill (Courtesy, greeting, smiling pitch and mood) menjadi faktor penting dalam Email Customer Service sekalipun dalam bentuk komunikasi tertulis. Hal ini dipantau secara detail dalam Riset CCSEI.\n'+
                      '2. Hard Skill adalah kemampuan Knowledge dan kemampuan memberikan solusi kepada pelanggan merupakan nilai tambah yang membedakan suatu Email Customer Service dengan yang lainnya. Disamping soal akurasi dalam berkomunikasi dalam menyampaikan solusi.'
                      ]]
        isi_tabel = [isi_tabel_enab, isi_tabel_enjoy1, isi_tabel_enjoy2]

    return isi_tabel



def FileTemplate():
    """
    Membuat template file PowerPoint kosong


    Return
    ------
    prs : objek template
    """
    prs = Presentation()
    return prs



def SubElement(parent, tagname, **kwargs):
    """
    Fungsi untuk mengedit kode XML python-pptx pada border tabel

    Parameter
    ---------
    parent : parent dari element
    tagname : nama tag untuk element
    **kwargs : parameter-parameter lainnya

    Return
    ------
    element : elemen xml yang sudah diganti
    """

    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def _set_cell_border(cell, border_color="000000", border_width='12700',
                    lnT=True, lnB=True, lnR=True, lnL=True):
    """
    Fungsi untuk mengatur border tabel.

    Parameter
    ---------
    border_width : tebal border tabel
    border_color : warna border tabel dengan kode hexadesimal
    lnT : border sebelah atas
    lnB : border sebelah bawah
    lnR : border sebelah kanan
    lnL : border sebelah kiri

    Return
    ------
    cell : cell tabel dengan border yang sudah di-edit
    """

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    if lnT==True:
        lnT = SubElement(tcPr, 'a:lnT', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(lnT, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(lnT, 'a:prstDash', val='solid')

    if lnB==True:
        lnB = SubElement(tcPr, 'a:lnB', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(lnB, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(lnB, 'a:prstDash', val='solid')

    if lnR==True:
        lnR = SubElement(tcPr, 'a:lnR', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(lnR, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(lnR, 'a:prstDash', val='solid')

    if lnL==True:
        lnL = SubElement(tcPr, 'a:lnL', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(lnL, 'a:solidFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(lnL, 'a:prstDash', val='solid')

    round_ = SubElement(lnB, 'a:round')
    headEnd = SubElement(lnB, 'a:headEnd', type='none', w='med', len='med')
    tailEnd = SubElement(lnB, 'a:tailEnd', type='none', w='med', len='med')
    return cell



def buat_tabel(template_ppt,nomor_slide,
               baris, kolom, kiri, atas, lebar, tinggi):
    """
    Membuat objek tabel di power point

    Parameter
    ---------
    template_ppt : template ppt yang dipakai
    nomor_slide : angka nomor slide (mulai dari 0)
    baris : jumlah baris tabel
    kolom : jumlah kolom tabel
    kiri : batas kiri tabel dengan tepi slide (dalam satian Inci)
    atas : batas atas tabel dengan tepi slide (dalam satian Inci)
    lebar : ukuran lebar seluruh tabel (dalam satian Inci)
    tinggi : ukuran tinggi seluruh tabel (dalam satian Inci)

    Return
    ------
    tabel : objek tabel di ppt
    """

    slide = template_ppt.slides[nomor_slide]
    shapes = slide.shapes

    rows = baris
    cols = kolom
    left = Inches(kiri)
    top = Inches(atas)
    width = Inches(lebar)
    height = Inches(tinggi)

    tabel = shapes.add_table(rows, cols, left, top, width, height).table

    return tabel


def halaman_slide(template_ppt):
    """
    Membuat halaman slide di tiap slide pada bagian kanan bawah

    Parameter
    ---------
    template_ppt : template ppt yang akan dibuat halaman slide-nya

    Return
    ------
    None
    """
    slide_ = template_ppt.slides
    for slide_halaman in slide_:
        if slide_.index(slide_halaman)!=0:

            left = Inches(8.25)
            top = Inches(7.1)
            width = Inches(1.5)
            height = Inches(0.25)
            txBox = slide_halaman.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "Halaman "+str(slide_.index(slide_halaman))+" dari "+str(len(slide_)-1)
            p.font.bold = True
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER


def slide_cover(template_ppt, kategori, tahun, list_bulan):
    """
    Membuat slide pada bagian cover

    Parameter
    ---------
    template_ppt : template ptt yang digunakan
    kategori : kategori SEI yang akan dibuat ('call center', 'online chat', 'twitter', 'email')
    tahun : tahun pembuatan laporan
    list_bulan : list array bulan yang ada di dalam dataframe

    Return
    ------
    None
    """

    slide = template_ppt.slide_layouts[6]
    slide = template_ppt.slides.add_slide(slide)

    #gambar halaman pertama
    img_path = "gambar/logo.png"
    left = Inches(.25)
    top = Inches(.25)
    width = Inches(2.5)
    height = Inches(2.25)
    slide.shapes.add_picture(img_path, left, top, width, height)

    #kotak
    left = Inches(3.25)
    width = Inches(6.5)
    height = Inches(2.25)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(204,105,255)
    shape.line.color.rgb = RGBColor(228,108,10)
    shape.line.width = Pt(3)
    shape.text_frame.paragraphs[0].text = ('LAPORAN AKHIR\n'+
                                         kategori.upper()+' SERVICE EXCELENCE INDEX '+str(tahun)+
                                         '\nPERIODE PEMANTAUAN\n'+
                                           str.upper(list_bulan[0])+' - '+str.upper(list_bulan[-1])+
                                          ' '+str(tahun))
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
    shape.text_frame.paragraphs[0].font.size = Pt(20)
    shape.text_frame.paragraphs[0].font.bold = True

    rows = 4
    cols = 3
    left = Inches(0.25)
    top = Inches(2.75)
    width = Inches(9)
    height = Inches(3)

    tb = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl =  tb._element.graphic.graphicData.tbl
    style_id = '{0505E3EF-67EA-436B-97B2-0124C06EBD24}' #style tabel MediumStyle4--Accent3 di ms office ppt
    tbl[0][-1].text = style_id
    tabel_cover = tb.table

    # set column widths
    tabel_cover.columns[0].width = Inches(2)
    tabel_cover.columns[1].width = Inches(.5)
    tabel_cover.columns[2].width = Inches(7)

    isi_tabel = [['No. Laporan',':',''],
                 ['Tipe Laporan',':',''],
                 ['Nama Client',':',''],
                ['Industri',':','']]

    # write column headings
    for i in range(len(isi_tabel)):
        tabel_cover.rows[i].height = Inches(0.5)
        for j in range(len(isi_tabel[i])):
            _set_cell_border(tabel_cover.cell(i,j))
            tabel_cover.cell(i,j).text = isi_tabel[i][j]
            tabel_cover.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
            tabel_cover.cell(i,j).text_frame.paragraphs[0].font.size = Pt(20)
            tabel_cover.cell(i,j).text_frame.paragraphs[0].font.bold = True
            tabel_cover.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
            tabel_cover.cell(i,j).fill.solid()
            if i==0 or i%2==0:
                tabel_cover.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                tabel_cover.cell(i,j).fill.fore_color.rgb = RGBColor(231, 231, 231)
            else:
                tabel_cover.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
                tabel_cover.cell(i,j).fill.fore_color.rgb = RGBColor(203, 203,203)
            if j==1:
                tabel_cover.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    #text box halaman pertama
    left = Inches(4.15)
    top = Inches(5)
    width = Inches(3)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Dilaporkan olah:"
    p.font.bold = True
    p.font.size = Pt(18)

    #gambar halaman pertama
    img_path = "gambar/halaman awal.png"
    left = Inches(3.75)
    top = Inches(5.25)
    width = Inches(2.75)
    height = Inches(1.25)
    slide.shapes.add_picture(img_path, left, top, width, height)


    #text box halaman pertama
    left = Inches(0.25)
    top = Inches(6.75)
    width = Inches(9.5)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = ("All rights reserved.  No part of this publication may be reproduced, stored in a retrieval system, or transmitted in any form \nor by any"+
              "means, electronic, mechanical, photocopying, recording, or otherwise, \nwithout prior written permission of The CARRE - Center for Customer "+
              "Satisfaction and Loyalty.")
    p.font.bold = True
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER



def footer(gambar_footer, slide):
    """
    Meletakkan gambar footer di bawah tengah

    Parameter
    ---------
    gambar_footer : nama file gambar yang akan dipakai
    slide : nomor slide yang akan diberi gambar footer

    Return
    ------
    None
    """

    footer_left = Inches(4.4)
    footer_top = Inches(7.1)
    footer_width = Inches(1.25)
    footer_height = Inches(0.35)
    footer_ = slide.shapes.add_picture(gambar_footer, footer_left, footer_top, footer_width, footer_height)



def slide_pembuka(template_ppt, gambar_footer):
    """
    Membuat slide pembuka setelah slide cover

    Parameter
    ---------
    template_ppt : template ptt yang digunakan
    gambar_footer : nama file gambar yang akan dipakai

    Return
    ------
    None
    """

    slide = template_ppt.slide_layouts[6]
    slide = template_ppt.slides.add_slide(slide)

    #text box halaman pertama
    left = Inches(3.2)
    top = Inches(2)
    width = Inches(3)
    height = Inches(3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Monitoring dikelola olah:"
    p.font.bold = True
    p.font.size = Pt(24)


    #gambar halaman pertama
    img_path = "gambar/halaman awal.png"
    left = Inches(2.25)
    top = Inches(2.5)
    width = Inches(5.5)
    height = Inches(2.65)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)

    footer(gambar_footer,slide)


def slide_daftar_isi(template_ppt, gambar_footer, kriteria):
    """
    Membuat slide berisi tabel daftar isi

    Parameter
    ---------
    template_ppt : template ptt yang digunakan
    gambar_footer : nama file gambar yang akan dipakai
    kriteria : kategori SEI yang akan dibuat ('call center', 'online chat', 'twitter', 'email')

    Return
    ------
    None
    """

    slide_daftar_isi = template_ppt.slide_layouts[6]
    slide = template_ppt.slides.add_slide(slide_daftar_isi)
    shapes = slide.shapes

    #text box halaman pertama
    left = Inches(0.5)
    top = Inches(0.25)
    width = Inches(9)
    height = Inches(3)
    txBox = shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "DAFTAR ISI"
    p.font.bold = True
    p.font.size = Pt(25)

    if kriteria=='email':
        isi_tabel = [['No','Bagian'],
                    ['1.','Frame Work ESEI'],
                    ['2.','Definisi Dimensi dan KPI ESEI'],
                    ['3.','Kinerja Dimensi Enabling'],
                    ['4.','Kinerja Dimensi Enjoying'],
                    ['5.','Engagement Index'],
                    ['6.','Area Perbaikan']]
    elif kriteria=='twitter':
        isi_tabel = [['No','Bagian'],
                    ['1.','Frame Work TSEI'],
                    ['2.','Definisi Dimensi dan KPI TSEI'],
                    ['3.','Kinerja Dimensi Engaging'],
                    ['4.','Kinerja Dimensi Human Touching'],
                    ['5.','Kinerja Dimensi Navigating'],
                    ['6.','Engagement Index'],
                    ['7.','Area Perbaikan']]
    elif kriteria=='online chat':
            isi_tabel = [['No','Bagian'],
                        ['1.','Frame Work OCSEI'],
                        ['2.','Definisi Dimensi dan KPI OCSEI'],
                        ['3.','Kinerja Dimensi Engaging'],
                        ['4.','Kinerja Dimensi Human Touching'],
                        ['5.','Kinerja Dimensi Navigating'],
                        ['6.','Engagement Index'],
                        ['7.','Area Perbaikan']]                    
    elif kriteria=='call center':
        isi_tabel = [['No','Bagian'],
                    ['1.','Frame Work CCSEI'],
                    ['2.','Definisi Dimensi dan KPI CCSEI'],
                    ['3.','Kinerja Dimensi Access'],
                    ['4.','Kinerja Dimensi System & Procedure'],
                    ['5.','Kinerja Dimensi People'],
                    ['6.','Engagement Index'],
                    ['7.','Area Perbaikan']]

    rows = len(isi_tabel)
    cols = len(isi_tabel[0])
    left = Inches(0.5)
    top = Inches(0.75)
    width = Inches(9)
    height = Inches(5)

    tabel_daftar_isi = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    tabel_daftar_isi.columns[0].width = Inches(1.)
    tabel_daftar_isi.columns[1].width = Inches(8)

    # write column headings
    for i in range(len(isi_tabel)):
        for j in range(len(isi_tabel[i])):
            _set_cell_border(tabel_daftar_isi.cell(i,j))
            tabel_daftar_isi.cell(i,j).text = isi_tabel[i][j]
            tabel_daftar_isi.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)

            if i==0:
                tabel_daftar_isi.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                tabel_daftar_isi.cell(i,j).text_frame.paragraphs[0].font.size = Pt(20)
                tabel_daftar_isi.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
                tabel_daftar_isi.cell(i,j).fill.solid()
                tabel_daftar_isi.cell(i,j).fill.fore_color.rgb = RGBColor(255, 128, 0)

            else:
                tabel_daftar_isi.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
                if i%2!=0:
                    tabel_daftar_isi.cell(i,j).fill.solid()
                    tabel_daftar_isi.cell(i,j).fill.fore_color.rgb = RGBColor(255, 255, 0)
                else:
                    tabel_daftar_isi.cell(i,j).fill.solid()
                    tabel_daftar_isi.cell(i,j).fill.fore_color.rgb = RGBColor(255, 255, 153)

                if j==0:
                    tabel_daftar_isi.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    footer(gambar_footer, slide)



def slide_transisi(template_ppt, gambar_footer, tulisan):
    """
    Membuat slide transisi dengan box dan tulisannya

    Parameter
    ---------
    template_ppt : template ptt yang digunakan
    gambar_footer : nama file gambar yang akan dipakai
    tulisan : tulisan di dalam text box

    Return
    ------
    None
    """
    slide_transisi = template_ppt.slide_layouts[6]
    slide = template_ppt.slides.add_slide(slide_transisi)
    shapes = slide.shapes

    #kotak
    left = Inches(0.5)
    top = Inches(2.75)
    width = Inches(9)
    height = Inches(2)
    shape = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(127,127,127)
    shape.line.color.rgb = RGBColor(0,0,0)

    #text
    if "Frame Work" in tulisan:
        top = Inches(3.125)
    else:
        top = Inches(3.4)
    #left = Inches(1)
    #width = Inches(9)
    #height = Inches(3)
    txBox = shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str.upper(tulisan)
    p.font.bold = True
    p.font.size = Pt(30)
    p.alignment = PP_ALIGN.CENTER

    footer(gambar_footer, slide)



def slide_framework(template_ppt, gambar_footer, kriteria):
    """
    Membuat slide ilustrasi framework SEI sesuai kategori SEI-nya

    Parameter
    ---------
    template_ppt : template ptt yang digunakan
    gambar_footer : nama file gambar yang akan dipakai
    kriteria : kategori SEI yang akan dibuat ('call center', 'online chat', 'twitter', 'email')

    Return
    ------
    None
    """

    slide_framework = template_ppt.slide_layouts[6]
    slide = template_ppt.slides.add_slide(slide_framework)
    shapes = slide.shapes

    if kriteria=='call center':
        teks = 'CCSEI'
        gambar = "framework "+kriteria+".png"
    elif kriteria =='twitter':
        teks = "TSEI"
        gambar = "framework "+kriteria+".jpeg"
    elif kriteria == 'online chat':
        teks = "OCSEI"
        gambar = "framework "+kriteria+".png"    
    elif kriteria=='email':
        teks = "ESEI"
        gambar = "framework "+kriteria+".jpeg"

    #judul
    left = Inches(0.5)
    top = Inches(0.25)
    width = Inches(9)
    height = Inches(3)
    txBox = shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "FRAME WORK "+teks
    p.font.bold = True
    p.font.size = Pt(25)

    #gambar
    img_path = "gambar/"+gambar
    left = Inches(0.5)
    top = Inches(0.85)
    width = Inches(9)
    height = Inches(5.75)
    pic = slide.shapes.add_picture(img_path, left, top, width, height)

    footer(gambar_footer, slide)



def slide_definisi_dimensi(template_ppt, gambar_footer, kriteria):
    """
    Membuat slide yang berisi definisi tiap dimensi yang dipakai dalam framework SEI

    Parameter
    ---------
    template_ppt : template ptt yang digunakan
    gambar_footer : nama file gambar yang akan dipakai
    kriteria : kategori SEI yang akan dibuat ('call center', 'online chat', 'twitter', 'email')

    Return
    ------
    None
    """

    isi_tabel = tabel_definisi_dimensi(kriteria)
    for i_slide in range(len(isi_tabel)):
        data_isi_tabel = isi_tabel[i_slide]

        slide_framework = template_ppt.slide_layouts[6]
        slide = template_ppt.slides.add_slide(slide_framework)
        shapes = slide.shapes

        if kriteria=='call center':
            teks = 'CCSEI'
        elif kriteria =='twitter' or kriteria == 'online chat':
            teks = "TCSEI"
        elif kriteria=='email':
            teks = "ECSEI"

        #judul
        left = Inches(0.5)
        top = Inches(0.25)
        width = Inches(9)
        height = Inches(3)
        txBox = shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "DEFINISI DIMENSI DAN KPI "+teks
        p.font.bold = True
        p.font.size = Pt(25)

        #tabel1
        rows = 2
        cols = 1
        left = Inches(0.5)
        top = Inches(0.75)
        width = Inches(9)
        height = Inches(1.5)
        tabel_definisi_dimensi1 = shapes.add_table(rows, cols, left, top, width, height).table

        # #tabel2
        if i_slide==0 or i_slide==1:
            top = Inches(2.15)
        else:
            top = Inches(2.0)
        tabel_definisi_dimensi2 = shapes.add_table(rows, cols, left, top, width, height).table

        if len(data_isi_tabel)==2:
            list_definisi_dimensi = [tabel_definisi_dimensi1,
                                   tabel_definisi_dimensi2]
        elif len(data_isi_tabel)==3:
            # #tabel3
            if i_slide==0:
                top = Inches(3.4)
            elif i_slide==1:
                top = Inches(3.55)
            else:
                top = Inches(3.25)
            tabel_definisi_dimensi3 = shapes.add_table(rows, cols, left, top, width, height).table
            list_definisi_dimensi = [tabel_definisi_dimensi1,
                                   tabel_definisi_dimensi2,
                                   tabel_definisi_dimensi3]
        else:
            # #tabel3
            if i_slide==0:
                top = Inches(3.4)
            elif i_slide==1:
                top = Inches(3.55)
            else:
                top = Inches(3.25)
            tabel_definisi_dimensi3 = shapes.add_table(rows, cols, left, top, width, height).table
            # #tabel4
            if i_slide==0:
                top = Inches(4.8)
            else:
                top = Inches(4.8)
            tabel_definisi_dimensi4 = shapes.add_table(rows, cols, left, top, width, height).table
            list_definisi_dimensi = [tabel_definisi_dimensi1,
                                           tabel_definisi_dimensi2,
                                           tabel_definisi_dimensi3,
                                           tabel_definisi_dimensi4]

        for i in range(len(list_definisi_dimensi)):
            tabel = list_definisi_dimensi[i]
            tabel.rows[0].height = Inches(0.25)
            for j in range(len(data_isi_tabel[i])):
                _set_cell_border(tabel.cell(j,0))
                tabel.cell(j,0).text = data_isi_tabel[i][j]
                tabel.cell(j,0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                tabel.cell(j,0).vertical_anchor = MSO_ANCHOR.MIDDLE
                tabel.cell(j,0).fill.solid()

                if i==0 and j==0:
                    tabel.cell(j,0).fill.fore_color.rgb = RGBColor(255, 128, 0)
                    tabel.cell(j,0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                elif i!=0 and j==0:
                    tabel.cell(j,0).fill.fore_color.rgb = RGBColor(255, 255, 0)
                    tabel.cell(j,0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                else:
                    tabel.cell(j,0).fill.fore_color.rgb = RGBColor(255, 255, 153)
                    for prg in range(len(tabel.cell(j,0).text_frame.paragraphs)):
                        tabel.cell(j,0).text_frame.paragraphs[prg].alignment = PP_ALIGN.JUSTIFY
                        tabel.cell(j,0).text_frame.paragraphs[prg].font.size = Pt(16)

        footer(gambar_footer, slide)



def plot_grafik_tracking(template_ppt, gambar_footer, list_bulan, data, dimensi, kpi, client):
    """
    Membuat slide grafik tracking kpi SEI

    Parameter
    ---------
    template_ppt : template ptt yang digunakan
    gambar_footer : nama file gambar yang akan dipakai
    list_bulan : list array bulan yang ada di dalam dataframe
    data : file data yang akan ditampilkan di slide
    dimensi : nama dimensi yang akan ditampilkan
    kpi : nama kpi yang akan ditampilkan
    client : nama client yang dibuatkan laporan

    Return
    ------
    None
    """

    slide_grafik = template_ppt.slide_layouts[6]
    slide = template_ppt.slides.add_slide(slide_grafik)

    #text
    left = Inches(2)
    top = Inches(0.25)
    width = Inches(6)
    height = Inches(2.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Grafik Kinerja Bulanan Dimensi "+dimensi.upper()+"\nKPI "+kpi+"\n"+client
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(16)

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4.75)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(155,187,89)
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(5)
    shape.shadow.inherit = False

    left = Inches(1.35)
    top = Inches(1.92)
    width = Inches(7.23)
    height = Inches(3.65)
    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(255,255,255)
    shape1.line.color.rgb = RGBColor(255,255,255)
    shape1.shadow.inherit = False

    #data
    data_kpi = pd.read_excel(data, sheet_name='kpi', engine='openpyxl')
    data_kpi['Dimensi'] = data_kpi['Dimensi'].fillna(method='ffill')
    data_kpi['KPI'] = data_kpi['KPI'].fillna(method='ffill')
    data_kpi = data_kpi.set_index(['Dimensi','KPI','Brand'])

    data_client = data_kpi.loc[(dimensi,kpi,client)][:len(list_bulan)]
    data_client = round(data_client,1)
    data_kosong = [None for i in range(13) if i>len(data_client)]
    data_client = data_client.append(pd.Series(data_kosong))
    data_client = data_client.values
    data_client = ['#N/A' if str(x)=='nan' else x for x in data_client]

    # define chart data ---------------------
    chart_data = CategoryChartData()
    chart_data.categories = nama_bulan[:len(list_bulan)]#data_client.index
    chart_data.add_series(client, data_client)

    x, y, cx, cy = Inches(0.75), Inches(1.75), Inches(8), Inches(4.25)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = False

    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 50
    chart.value_axis.major_unit = 10
    chart.value_axis.tick_labels.number_format = '#"%"'
    chart.value_axis.tick_labels.font.size = Pt(12)

    chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    chart.category_axis.tick_labels.font.size = Pt(12)

    chart.plots[0].has_data_label = True
    chart.plots[0].series[0].format.line.color.rgb = RGBColor(255, 0, 0)
    chart.plots[0].series[0].format.line.width = Pt(3)

    for i in range(len(list_bulan)):
        label = data_client[i]
        chart.plots[0].series[0].points[i].data_label.text_frame.text = str(label)+"%"
        chart.plots[0].series[0].points[i].data_label.ShowBubbleSize = True
        chart.plots[0].series[0].points[i].data_label.font.size = Pt(12)
        chart.plots[0].series[0].marker.style = XL_MARKER_STYLE.CIRCLE
        chart.plots[0].series[0].marker.format.fill.solid()
        chart.plots[0].series[0].marker.format.fill.fore_color.rgb = RGBColor(0,0,0)
        if type(label)!=str:
            if label > 60:
                chart.plots[0].series[0].points[i].data_label.position = XL_LABEL_POSITION.BELOW
            else:
                chart.plots[0].series[0].points[i].data_label.position = XL_LABEL_POSITION.ABOVE
        else:
            chart.plots[0].series[0].points[i].data_label.position = XL_LABEL_POSITION.ABOVE

    footer(gambar_footer, slide)



def plot_tabel_data_semester(template_ppt, gambar_footer, list_bulan, data_xls, dimensi_, subkpi_,
                             client, kriteria=None, bobot=None):
    """
    Membuat slide tabel data per subkpi pada semester 1 dan semester 2

    Parameter
    ---------
    template_ppt : template ptt yang digunakan
    gambar_footer : nama file gambar yang akan dipakai
    list_bulan : list array bulan yang ada di dalam dataframe
    data_xls : file data yang akan ditampilkan di slide
    dimensi_ : nama dimensi yang akan ditampilkan
    subkpi_ : nama subkpi_ yang akan ditampilkan
    client : nama client yang dibuatkan laporan
    kriteria : kategori SEI yang akan dibuat ('call center', 'online chat', 'twitter', 'email')
    bobot : dataframe yang berisi daftar bobot perhitungan indeks

    Return
    ------
    None
    """

    jumlah_semester = 2

    #data
    data_mentah = pd.read_excel(data_xls, sheet_name='data mentah', engine='openpyxl')
    
    data_aspek = pd.read_excel(data_xls, sheet_name='aspek', engine='openpyxl')
    data_aspek['Dimensi'] = data_aspek['Dimensi'].fillna(method='ffill')
    data_aspek['KPI'] = data_aspek['KPI'].fillna(method='ffill')
    data_aspek['SUB KPI'] = data_aspek['SUB KPI'].fillna(method='ffill')
    data_aspek['Nomor'] = data_aspek['Nomor'].fillna(method='ffill')
    data_aspek['Kode'] = data_aspek['Kode'].fillna(method='ffill')
    data_aspek['Aspek'] = data_aspek['Aspek'].fillna(method='ffill')
    data_aspek = data_aspek.set_index(['Dimensi','KPI','SUB KPI','Nomor','Kode','Aspek','Brand'])

    df_subkpi = data_aspek[data_aspek.index.get_level_values('SUB KPI').isin([subkpi_])]
    
    list_aspek = df_subkpi.index.get_level_values('Aspek').unique().values
    
    data_subkpi = pd.read_excel(data_xls, sheet_name='subkpi', index_col=[0,1,2,3], engine='openpyxl')
    data_subkpi = data_subkpi[data_subkpi.index.get_level_values("SUB KPI").isin([subkpi_])]

    if (kriteria=='twitter') or (kriteria=='online chat') or (kriteria=='email'):
        crosstab_rata_transpose = pd.read_excel(data_xls, sheet_name = 'crosstab transpose',
                                               header=[0,1],index_col=[0], engine='openpyxl')
        crosstab_rata_transpose = crosstab_rata_transpose.reset_index()
        crosstab_rata_transpose = crosstab_rata_transpose.rename(columns = {'KODE':'Aspek'})

    if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility'):
        var= bobot[['ASPEK','SUB ASPEK']][:4].copy()
        asp = df_subkpi.index.get_level_values('Aspek').unique().to_list()
        aspek_feas = []
        for i in range(len(asp)):
            if i==0:
                var['ASPEK1'] = var['ASPEK']+var['SUB ASPEK']
                var['ASPEK1']=var['ASPEK1'].str.replace('Ada','')
                var['ASPEK1']=var['ASPEK1'].str.replace('Terlihat','')
                var['ASPEK1']=var['ASPEK1'].str.replace('Di',' di')

                apa = crosstab_rata_transpose[:4].copy()                
                apa['Aspek']=var['ASPEK1']
                apa = apa.set_index('Aspek')
                apa = apa.stack()
                apa = apa.reindex(columns=nama_bulan[:len(list_bulan)])
                apa['~Average'] = round(apa.mean(axis=1),1)

                list_aspek = apa[apa.index.get_level_values('Brand').isin([client])].index.values
                list_aspek = [x[0] for x in list_aspek]
            else:
                list_aspek = [asp[i]]
            aspek_feas.append(list_aspek)

    elif (kriteria=='email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')):
        var= bobot[['ASPEK','SUB ASPEK']].copy()
        aspek_feas = []
        for i in range(3):
            asp = df_subkpi.index.get_level_values('Aspek').unique().to_list()
            if i==0:
                if subkpi_=='Service Standard Consistency':
                    asp = [x for x in asp if ('pembuka' in x)]
                elif subkpi_ == 'Soft Skill':
                    asp = [x for x in asp if ('etika yang baik' in x)]
                subaspek_ = var[var['ASPEK'].isin(asp)]['SUB ASPEK'].to_list()
                idx = var[var['ASPEK'].isin(asp)].index.to_list()
                apa1 = crosstab_rata_transpose.loc[idx, :].copy()
                apa1 = apa1.reset_index(drop=True)
                apa1['Aspek'] = pd.Series(subaspek_)
                apa1 = apa1.set_index('Aspek')
                apa1 = apa1.stack()
                apa1 = apa1.reindex(columns=nama_bulan[:len(list_bulan)])
                apa1['~Average'] = round(apa1.mean(axis=1),1)

                list_aspek = apa1[apa1.index.get_level_values('Brand').isin([client])].index.values
                list_aspek = [x[0] for x in list_aspek]
            elif i==1:
                if subkpi_=='Service Standard Consistency':
                    asp = [x for x in asp if ('penutup' in x)]
                elif subkpi_=='Soft Skill':
                    asp = [x for x in asp if ('Service yang baik' in x)]
                subaspek_ = var[var['ASPEK'].isin(asp)]['SUB ASPEK'].to_list()
                idx = var[var['ASPEK'].isin(asp)].index.to_list()
                apa2 = crosstab_rata_transpose.loc[idx, :].copy()
                apa2 = apa2.reset_index(drop=True)
                apa2['Aspek'] = pd.Series(subaspek_)
                apa2 = apa2.set_index('Aspek')
                apa2.columns = apa2.columns.set_levels(nomor_bulan, level=0)
                apa2 = apa2.stack(dropna=False)
                apa2.columns = nama_bulan[:len(list_bulan)]
                apa2['~Average'] = round(apa2.mean(axis=1),1)

                list_aspek = apa2[apa2.index.get_level_values('Brand').isin([client])].index.values
                list_aspek = [x[0] for x in list_aspek]
            else:
                list_aspek = df_subkpi.index.get_level_values('Aspek').unique().to_list()

            aspek_feas.append(list_aspek)

    else:
        aspek_feas = [1]

    for asp in range(len(aspek_feas)):
        for semester in range(jumlah_semester):
            slide_tabel_nilai = template_ppt.slide_layouts[6]
            slide = template_ppt.slides.add_slide(slide_tabel_nilai)
            shapes = slide.shapes

            #judul
            left = Inches(0.5)
            top = Inches(0.25)
            width = Inches(9)
            height = Inches(3)
            txBox = shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "Tabel Kinerja Bulanan KPI dan Atribut Dalam\nDimensi "+dimensi_.upper()+"\n"+client
            p.font.bold = True
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER

            cols = 8
            if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility') and (asp==1):
                rows = len([aspek_feas[asp]])+3
            elif (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility') and (asp==0):
                rows = len(aspek_feas[asp])+3
            elif (kriteria=='email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')):
                rows = len(aspek_feas[asp])+3
            else:
                rows = len(list_aspek)+3
            top = Inches(1.5)
            height = Inches(1.5)
            left = Inches(0.25)
            width = Inches(9.25)

            tabel_plot = shapes.add_table(rows, cols, left, top, width, height).table
            tabel_plot.columns[0].width = Inches(.25)
            tabel_plot.columns[1].width = Inches(3.25)
            lebar_kolom_lainnya = Inches(1)
            for ke in range(cols-2):
                tabel_plot.columns[ke+2].width = lebar_kolom_lainnya

            list_all = []
            for j in range(rows):
                list_temp = []
                if j==0:
                    if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_=='Feasibility'):
                        header1 = ['','ATRIBUT %s %s'%(subkpi_.upper(), str(asp+1)),'Januari','Februari','Maret','April','Mei','Juni']
                        header2 = ['','ATRIBUT %s %s'%(subkpi_.upper(), str(asp+1)),'Juli','Agustus','September','Oktober','November','Desember']
                    elif (kriteria=='email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')):
                        if asp==0:
                            if subkpi_=='Service Standard Consistency':
                                header1 = ['','ATRIBUT %s - %s'%(subkpi_.upper(), 'BAGIAN PEMBUKA'),'Januari','Februari','Maret','April','Mei','Juni']
                                header2 = ['','ATRIBUT %s - %s'%(subkpi_.upper(), 'BAGIAN PEMBUKA'),'Juli','Agustus','September','Oktober','November','Desember']
                            elif subkpi_=='Soft Skill':
                                header1 = ['','ATRIBUT %s - %s'%(subkpi_.upper(), 'ETIKA YANG BAIK'),'Januari','Februari','Maret','April','Mei','Juni']
                                header2 = ['','ATRIBUT %s - %s'%(subkpi_.upper(), 'ETIKA YANG BAIK'),'Juli','Agustus','September','Oktober','November','Desember']
                        elif asp==1:
                            if subkpi_=='Service Standard Consistency':
                                header1 = ['','ATRIBUT %s - %s'%(subkpi_.upper(), 'BAGIAN PENUTUP'),'Januari','Februari','Maret','April','Mei','Juni']
                                header2 = ['','ATRIBUT %s - %s'%(subkpi_.upper(), 'BAGIAN PENUTUP'),'Juli','Agustus','September','Oktober','November','Desember']
                            elif subkpi_=='Soft Skill':
                                header1 = ['','ATRIBUT %s - %s'%(subkpi_.upper(), 'BAHASA SERVICE YANG BAIK'),'Januari','Februari','Maret','April','Mei','Juni']
                                header2 = ['','ATRIBUT %s - %s'%(subkpi_.upper(), 'BAHASA SERVICE YANG BAIK'),'Juli','Agustus','September','Oktober','November','Desember']
                        else:
                            header1 = ['','ATRIBUT %s'%(subkpi_.upper()),'Januari','Februari','Maret','April','Mei','Juni']
                            header2 = ['','ATRIBUT %s'%(subkpi_.upper()),'Juli','Agustus','September','Oktober','November','Desember']
                    else:
                        header1 = ['','ATRIBUT '+subkpi_.upper(),'Januari','Februari','Maret','April','Mei','Juni']
                        header2 = ['','ATRIBUT '+subkpi_.upper(),'Juli','Agustus','September','Oktober','November','Desember']
                    if semester==0:
                        header = header1
                        list_temp = header
                    else:
                        header = header2
                        list_temp = header

                    for k in range(cols):
                        _set_cell_border(tabel_plot.cell(j,k))
                        tabel_plot.cell(j,k).text = list_temp[k]
                        tabel_plot.cell(j,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                        tabel_plot.cell(j,k).text_frame.paragraphs[0].font.size = Pt(11)
                        tabel_plot.cell(j,k).vertical_anchor = MSO_ANCHOR.MIDDLE
                        tabel_plot.cell(j,k).fill.solid()
                        tabel_plot.cell(j,k).fill.fore_color.rgb = RGBColor(155,187,89)
                        tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    tabel_plot.cell(j,0).merge(tabel_plot.cell(j,1))

                elif j<(rows-3)+1:
                    if (subkpi_ == 'Feasibility') and (asp==0) and (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email'):
                        value_client = apa[apa.index.get_level_values('Brand').isin([client])].values[j-1]
                    elif (subkpi_ == 'Feasibility') and (asp==1) and (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email'):
                        temp_df_asp = df_subkpi[df_subkpi.index.get_level_values('SUB KPI').isin([subkpi_])]
                        value_client = temp_df_asp[temp_df_asp.index.get_level_values('Aspek').isin(aspek_feas[asp])]
                        #display(value_client[value_client.index.get_level_values('Brand').isin([client])])
                        value_client = value_client[value_client.index.get_level_values('Brand').isin([client])].values[0]

                    elif ((subkpi_=='Service Standard Consistency') or (subkpi_=='Soft Skill')) and (kriteria=='email') and (asp!=len(aspek_feas)-1):
                        if asp==0:
                            value_client = apa1[apa1.index.get_level_values('Brand').isin([client])].values[j-1]
                        else:
                            value_client = apa2[apa2.index.get_level_values('Brand').isin([client])].values[j-1]
                    elif ((subkpi_=='Service Standard Consistency') or (subkpi_=='Soft Skill')) and (kriteria=='email') and (asp==len(aspek_feas)-1):
                        temp_df_asp = df_subkpi[df_subkpi.index.get_level_values('SUB KPI').isin([subkpi_])]
                        value_client = temp_df_asp[temp_df_asp.index.get_level_values('Aspek').isin(aspek_feas[asp])]
                        value_client = value_client[value_client.index.get_level_values('Brand').isin([client])].values[0]

                    else:
                        temp_df_asp = df_subkpi[df_subkpi.index.get_level_values('Aspek').isin([list_aspek[j-1]])]
                        value_client = temp_df_asp[temp_df_asp.index.get_level_values('Brand').isin([client])].values[0]

                    #jika data yang ada maksimal hanya terisi sampai bulan Juni, pakai data ini
                    if len(list_bulan)>0 and len(list_bulan)<=6:
                        if semester==0:
                            list_value = value_client[:len(list_bulan)]
                        else:
                            list_value = ['']*6

                    #jika data yang ada terisi sampai melewati bulan Juli, pakai data ini
                    else:
                        if semester==0:
                            list_value = value_client[:len(list_bulan)]
                        else:
                            list_value = value_client[6:6+len(nama_bulan)-len(list_bulan)]

                    list_temp.extend([j])

                    if (subkpi_ == 'Feasibility') and (asp==0) and (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email'):
                        list_temp.extend([aspek_feas[asp][j-1]])
                    elif (subkpi_ == 'Feasibility') and(asp==1) and (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email'):
                        list_temp.extend(aspek_feas[asp])
                    elif ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')) and (kriteria=='email'):
                        list_temp.extend([aspek_feas[asp][j-1]])
                    else:
                        list_temp.extend([list_aspek[j-1]])

                    list_temp.extend(list_value)

                    list_all.append(list_temp)

                    if len(list_temp)!=8:
                        kurang = 8-len(list_temp)
                        for kurangnya in range(kurang):
                            list_temp.extend([''])

                    for k in range(cols):
                        t_value = list_temp[k]
                        if str(t_value)=='nan':
                            val_ = '#N/A'
                        elif type(t_value)==np.float64:
                            val_ = str(round(t_value,1))+"%"
                        else:
                            val_ = str(t_value)
                        _set_cell_border(tabel_plot.cell(j,k))
                        tabel_plot.cell(j,k).text = val_
                        tabel_plot.cell(j,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                        tabel_plot.cell(j,k).text_frame.paragraphs[0].font.size = Pt(10)

                        if k==1:
                            tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                        else:
                            tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            tabel_plot.cell(j,k).vertical_anchor = MSO_ANCHOR.MIDDLE

                        tabel_plot.cell(j,k).fill.solid()
                        if j%2!=0:
                            tabel_plot.cell(j,k).fill.fore_color.rgb = RGBColor(255, 128, 0)
                        else:
                            tabel_plot.cell(j,k).fill.fore_color.rgb = RGBColor(255, 255, 153)


                elif j==(rows-3)+1:
                    list_temp.extend([''])
                    list_temp.extend(['TOTAL SCORE'])

                    # menghitung total khusus untuk data maksimal terisi di bawah bulan Juni
                    if len(list_bulan)>0 and len(list_bulan)<=6:
                        if semester==0:
                            for k in list_bulan:
                                if ('Feasibility'==subkpi_) and (asp==0) and (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email'):
                                    rata_ = apa[apa.index.get_level_values('Brand').isin([client])].values[0][nama_bulan.index(k)]
                                    rata_ = round(rata_,1)
                                    #display("feasibility 1", apa)
                                elif ('Feasibility'==subkpi_) and (asp==1) and (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email'):
                                    rata_ =value_client[-1]
                                    rata_ = round(rata_,1)
                                    #display("feasibility 2", value_client)

                                elif ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')) and (kriteria=='email') and (asp!=len(aspek_feas)-1):
                                    if asp==0:
                                        rata_ =apa1[apa1.index.get_level_values('Brand').isin([client])].values[0][nama_bulan.index(k)]
                                        #display("SSC 0", apa1)
                                    else:
                                        rata_ =apa1[apa1.index.get_level_values('Brand').isin([client])].values[0][nama_bulan.index(k)]
                                        #display("SSC 1", apa1)
                                    rata_ = round(rata_,1)
                                elif ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')) and (kriteria=='email') and (asp==len(aspek_feas)-1):
                                    rata_ =value_client[-1]
                                    rata_ = round(rata_,1)
                                    #display("SSC 2", data_subkpi)

                                else:
                                    id_value = header.index(k)
                                    rata_skor = 0
                                    for l in range(len(list_all)):
                                        rata_skor += list_all[l][id_value]
                                    rata_ = round(rata_skor / len(list_aspek),1)
                                list_temp.extend([rata_])

                            if len(list_temp)!=8:
                                kurang = 8-len(list_temp)
                                for kurangnya in range(kurang):
                                    list_temp.extend([''])

                            for k in range(cols):
                                t_value = list_temp[k]
                                if str(t_value)=='nan':
                                    val_ = '#N/A'
                                elif type(t_value)==np.float64:
                                    val_ = str(t_value)+"%"
                                else:
                                    val_ = str(t_value)
                                _set_cell_border(tabel_plot.cell(j,k))
                                tabel_plot.cell(j,k).text = val_
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].font.size = Pt(10)
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].font.bold = True
                                tabel_plot.cell(j,k).vertical_anchor = MSO_ANCHOR.MIDDLE
                                tabel_plot.cell(j,k).fill.solid()
                                tabel_plot.cell(j,k).fill.fore_color.rgb = RGBColor(255, 0, 0)
                                if k >=2 :
                                    tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                                else:
                                    tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                        else:
                            for k in range(cols):
                                if k < 2:
                                    tabel_plot.cell(j,k).text = list_temp[k]
                                else:
                                    tabel_plot.cell(j,k).text = ''
                                _set_cell_border(tabel_plot.cell(j,k))
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].font.size = Pt(10)
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].font.bold = True
                                tabel_plot.cell(j,k).vertical_anchor = MSO_ANCHOR.MIDDLE
                                tabel_plot.cell(j,k).fill.solid()
                                tabel_plot.cell(j,k).fill.fore_color.rgb = RGBColor(255, 0, 0)
                                if k >=2 :
                                    tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                                else:
                                    tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    #                     tabel_plot.cell(j,0).merge(tabel_plot.cell(j,1))

                    #menghitung total jika data terisi melewati bulan Juli
                    else:
                        if semester==0:
                            nm_bln = nama_bulan[:6]
                        else:
                            nm_bln = [i for i in nama_bulan[6:] if i in list_bulan]

#                        if semester==0:
                        for k in nm_bln:
                            if ('Feasibility'==subkpi_) and (asp==0) and (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email'):
                                rata_ = apa[apa.index.get_level_values('Brand').isin([client])].values[0][nama_bulan.index(k)]
                                rata_ = round(rata_,1)
                                #display("feasibility 1", apa)
                            elif ('Feasibility'==subkpi_) and (asp==1) and (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email'):
                                rata_ =value_client[-1]
                                rata_ = round(rata_,1)
                                #display("feasibility 2", value_client)

                            elif ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')) and (kriteria=='email') and (asp!=len(aspek_feas)-1):
                                if asp==0:
                                    rata_ =apa1[apa1.index.get_level_values('Brand').isin([client])].values[0][nama_bulan.index(k)]
                                    #display("SSC 0", apa1)
                                else:
                                    rata_ =apa1[apa1.index.get_level_values('Brand').isin([client])].values[0][nama_bulan.index(k)]
                                    #display("SSC 1", apa1)
                                rata_ = round(rata_,1)
                            elif ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')) and (kriteria=='email') and (asp==len(aspek_feas)-1):
                                rata_ =value_client[-1]
                                rata_ = round(rata_,1)
                                #display("SSC 2", data_subkpi)

                            else:
                                id_value = header.index(k)
                                rata_skor = 0
                                for l in range(len(list_all)):
                                    rata_skor += list_all[l][id_value]
                                rata_ = round(rata_skor / len(list_aspek),1)
                            list_temp.extend([rata_])

                        if len(list_temp)!=8:
                            kurang = 8-len(list_temp)
                            for kurangnya in range(kurang):
                                list_temp.extend([''])

                        for k in range(cols):
                            t_value = list_temp[k]
                            if str(t_value)=='nan':
                                val_ = '#N/A'
                            elif type(t_value)==np.float64:
                                val_ = str(t_value)+"%"
                            else:
                                val_ = str(t_value)
                            _set_cell_border(tabel_plot.cell(j,k))
                            tabel_plot.cell(j,k).text = val_
                            tabel_plot.cell(j,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
                            tabel_plot.cell(j,k).text_frame.paragraphs[0].font.size = Pt(10)
                            tabel_plot.cell(j,k).text_frame.paragraphs[0].font.bold = True
                            tabel_plot.cell(j,k).vertical_anchor = MSO_ANCHOR.MIDDLE
                            tabel_plot.cell(j,k).fill.solid()
                            tabel_plot.cell(j,k).fill.fore_color.rgb = RGBColor(255, 0, 0)
                            if k >=2 :
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            else:
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

                else:
                    list_temp.extend([''])
                    list_temp.extend(['n Sample'])

                    #menghitung n sample jika data terisi maksimal di bulan Juni
                    if len(list_bulan)>0 and len(list_bulan)<=6:
                        if semester==0:
                            for k in list_bulan:
                                list_temp.extend([data_mentah.Periode.nunique()])

                            if len(list_temp)!=8:
                                kurang = 8-len(list_temp)
                                for kurangnya in range(kurang):
                                    list_temp.extend([''])

                            for k in range(cols):
                                t_value = list_temp[k]
                                if type(t_value)==np.float64:
                                    val_ = str(t_value)+"%"
                                else:
                                    val_ = str(t_value)
                                _set_cell_border(tabel_plot.cell(j,k))
                                tabel_plot.cell(j,k).text = val_
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].font.size = Pt(10)
                                tabel_plot.cell(j,k).vertical_anchor = MSO_ANCHOR.MIDDLE
                                tabel_plot.cell(j,k).fill.solid()
                                tabel_plot.cell(j,k).fill.fore_color.rgb = RGBColor(255, 255, 255)
                                if k >=2 :
                                    tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                                else:
                                    tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

                        else:
                            for k in range(cols):
                                if k < 2:
                                    tabel_plot.cell(j,k).text = list_temp[k]
                                else:
                                    tabel_plot.cell(j,k).text = ''
                                _set_cell_border(tabel_plot.cell(j,k))
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].font.size = Pt(10)
                                tabel_plot.cell(j,k).vertical_anchor = MSO_ANCHOR.MIDDLE
                                tabel_plot.cell(j,k).fill.solid()
                                tabel_plot.cell(j,k).fill.fore_color.rgb = RGBColor(255, 255, 255)
                                if k >=2 :
                                    tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                                else:
                                    tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    #                     tabel_plot.cell(j,0).merge(tabel_plot.cell(j,1))

                    #menghitung n sample jika data terisi melewati bulan Juli
                    else:
                        if semester==0:
                            nm_bln=nama_bulan[:6]
                        else:
                            nm_bln = [i for i in nama_bulan[6:] if i in list_bulan]

                        for k in nm_bln:
                            list_temp.extend([data_mentah.Periode.nunique()])

                        if len(list_temp)!=8:
                            kurang = 8-len(list_temp)
                            for kurangnya in range(kurang):
                                list_temp.extend([''])

                        for k in range(cols):
                            t_value = list_temp[k]
                            if type(t_value)==np.float64:
                                val_ = str(t_value)+"%"
                            else:
                                val_ = str(t_value)
                            _set_cell_border(tabel_plot.cell(j,k))
                            tabel_plot.cell(j,k).text = val_
                            tabel_plot.cell(j,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                            tabel_plot.cell(j,k).text_frame.paragraphs[0].font.size = Pt(10)
                            tabel_plot.cell(j,k).vertical_anchor = MSO_ANCHOR.MIDDLE
                            tabel_plot.cell(j,k).fill.solid()
                            tabel_plot.cell(j,k).fill.fore_color.rgb = RGBColor(255, 255, 255)
                            if k >=2 :
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            else:
                                tabel_plot.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    #                     tabel_plot.cell(j,0).merge(tabel_plot.cell(j,1))

            footer(gambar_footer, slide)


def plot_barchart_tabel(template_ppt, gambar_footer, list_bulan, data_xls, dimensi_, subkpi_,client, data_tahun,
                        kriteria=None, bobot=None):
    """
    Membuat slide barchart beserta tabelnya

    Parameter
    ---------
    template_ppt : template ptt yang digunakan
    gambar_footer : nama file gambar yang akan dipakai
    list_bulan : list array bulan yang ada di dalam dataframe
    data_xls : file data yang akan ditampilkan di slide
    dimensi_ : nama dimensi yang akan ditampilkan
    subkpi_ : nama subkpi_ yang akan ditampilkan
    client : nama client yang dibuatkan laporan
    data_tahun : tahun pembuatan laporan
    kriteria : kategori SEI yang akan dibuat ('call center', 'online chat', 'twitter', 'email')
    bobot : dataframe yang berisi daftar bobot perhitungan indeks

    Return
    ------
    None
    """

    #data
    data_mentah = pd.read_excel(data_xls, sheet_name='data mentah', engine='openpyxl')
    list_brand = data_mentah['Brand'].unique()

    data_aspek = pd.read_excel(data_xls, sheet_name='aspek', engine='openpyxl')
    data_aspek['Dimensi'] = data_aspek['Dimensi'].fillna(method='ffill')
    data_aspek['KPI'] = data_aspek['KPI'].fillna(method='ffill')
    data_aspek['SUB KPI'] = data_aspek['SUB KPI'].fillna(method='ffill')
    data_aspek['Nomor'] = data_aspek['Nomor'].fillna(method='ffill')
    data_aspek['Kode'] = data_aspek['Kode'].fillna(method='ffill')
    data_aspek['Aspek'] = data_aspek['Aspek'].fillna(method='ffill')
    data_aspek = data_aspek.set_index(['Dimensi','KPI','SUB KPI','Nomor','Kode','Aspek','Brand'])

    data_subkpi = pd.read_excel(data_xls, sheet_name='subkpi', engine='openpyxl')
    data_subkpi['Dimensi'] = data_subkpi['Dimensi'].fillna(method='ffill')
    data_subkpi['KPI'] = data_subkpi['KPI'].fillna(method='ffill')
    data_subkpi['SUB KPI'] = data_subkpi['SUB KPI'].fillna(method='ffill')
    data_subkpi = data_subkpi.set_index(['Dimensi','KPI','SUB KPI','Brand'])
    temp_sub = data_subkpi[data_subkpi.index.get_level_values('SUB KPI').isin([subkpi_])]

    temp_df_ = data_aspek[data_aspek.index.get_level_values('SUB KPI').isin([subkpi_])]
    list_aspek = temp_df_.index.get_level_values('Aspek').unique().values
    
    subkpi_ivr = od.subkpi_ivr

    if (kriteria=='twitter') or (kriteria=='online chat') or (kriteria=='email'):
        crosstab_rata_transpose = pd.read_excel(data_xls, sheet_name = 'crosstab transpose',
                                               header=[0,1],index_col=[0], engine='openpyxl')
        crosstab_rata_transpose = crosstab_rata_transpose.reset_index()
        crosstab_rata_transpose = crosstab_rata_transpose.rename(columns = {'KODE':'Aspek'})

    if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility'):
        asp = temp_df_.index.get_level_values('Aspek').unique().to_list()
        aspek_feas = []
        for i in range(len(asp)):
            if i==0:
                var= bobot[['ASPEK','SUB ASPEK']][:4].copy()
                var['ASPEK1'] = var['ASPEK']+var['SUB ASPEK']
                var['ASPEK1']=var['ASPEK1'].str.replace('Ada','')
                var['ASPEK1']=var['ASPEK1'].str.replace('Terlihat','')
                var['ASPEK1']=var['ASPEK1'].str.replace('Di',' di')

                apa = crosstab_rata_transpose[:4].copy()
                apa['Aspek']=var['ASPEK1']
                apa = apa.set_index('Aspek')
                apa.columns = apa.columns.set_levels(nomor_bulan, level=0)
                apa = apa.stack(dropna=False)
                apa.columns = nama_bulan[:len(list_bulan)]
                apa['~Average'] = round(apa.mean(axis=1),1)

                list_aspek = apa[apa.index.get_level_values('Brand').isin([client])].index.values
                list_aspek = [x[0] for x in list_aspek]
            else:
                list_aspek = [asp[i]]
            aspek_feas.append(list_aspek)

    elif (kriteria=='email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_ == 'Soft Skill')):
        var= bobot[['ASPEK','SUB ASPEK']].copy()
        aspek_feas = []
        for i in range(3):
            asp = temp_df_.index.get_level_values('Aspek').unique().to_list()
            if i==0:
                if subkpi_ == 'Service Standard Consistency':
                    asp = [x for x in asp if ('pembuka' in x)]
                elif subkpi_ == 'Soft Skill':
                    asp = [x for x in asp if ('etika yang baik' in x)]
                subaspek_ = var[var['ASPEK'].isin(asp)]['SUB ASPEK'].to_list()
                idx = var[var['ASPEK'].isin(asp)].index.to_list()
                apa1 = crosstab_rata_transpose.loc[idx, :].copy()
                apa1 = apa1.reset_index(drop=True)
                apa1['Aspek'] = pd.Series(subaspek_)
                apa1 = apa1.set_index('Aspek')
                apa1.columns = apa1.columns.set_levels(nomor_bulan, level=0)
                apa1 = apa1.stack(dropna=False)
                apa1.columns = nama_bulan[:len(list_bulan)]
                apa1['~Average'] = round(apa1.mean(axis=1),1)

                list_aspek = apa1[apa1.index.get_level_values('Brand').isin([client])].index.values
                list_aspek = [x[0] for x in list_aspek]
            elif i==1:
                if subkpi_=='Service Standard Consistency':
                    asp = [x for x in asp if ('penutup' in x)]
                elif subkpi_ =='Soft Skill':
                    asp = [x for x in asp if ('Service yang baik' in x)]
                subaspek_ = var[var['ASPEK'].isin(asp)]['SUB ASPEK'].to_list()
                idx = var[var['ASPEK'].isin(asp)].index.to_list()
                apa2 = crosstab_rata_transpose.loc[idx, :].copy()
                apa2 = apa2.reset_index(drop=True)
                apa2['Aspek'] = pd.Series(subaspek_)
                apa2 = apa2.set_index('Aspek')
                apa2.columns = apa2.columns.set_levels(nomor_bulan, level=0)
                apa2 = apa2.stack(dropna=False)
                apa2.columns = nama_bulan[:len(list_bulan)]
                apa2['~Average'] = round(apa2.mean(axis=1),1)

                list_aspek = apa2[apa2.index.get_level_values('Brand').isin([client])].index.values
                list_aspek = [x[0] for x in list_aspek]
            else:
                list_aspek = temp_df_.index.get_level_values('Aspek').unique().to_list()

            aspek_feas.append(list_aspek)

    else:
        aspek_feas = [1]

    for asp in range(len(aspek_feas)):
        list_aspek_sub = []
        if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_=='Feasibility'):
            list_aspek = aspek_feas[asp]
        elif (kriteria=='email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_ == 'Soft Skill')):
            list_aspek = aspek_feas[asp]

        if len(list_aspek)>4:
            temp_list_aspek = list_aspek[:4]
            list_aspek_sub.append(temp_list_aspek)
            if len(list_aspek[4:])>4:
                temp_list_aspek1 = list_aspek[4:8]
                temp_list_aspek2 = list_aspek[8:]
                list_aspek_sub.append(temp_list_aspek1)
                list_aspek_sub.append(temp_list_aspek2)
            else:
                temp_list_aspek1 = list_aspek[4:]
                list_aspek_sub.append(temp_list_aspek1)
        else:
            list_aspek_sub.append(list_aspek)

        for laspek in range(len(list_aspek_sub)):
            l_aspek_sub = list_aspek_sub[laspek]

            slide_barchart1 = template_ppt.slide_layouts[6]
            slide = template_ppt.slides.add_slide(slide_barchart1)
            shapes = slide.shapes

            # define chart data ---------------------
            chart_data = CategoryChartData()

            if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility'):
                l_asp = temp_df_.index.get_level_values('Aspek').unique().to_list()
                chart_kategori = list(temp_df_[temp_df_.index.get_level_values('Aspek').isin([l_asp[asp]])].index.get_level_values('Brand').values)
                chart_kategori = chart_kategori[::-1]
                for ktg in range(len(chart_kategori)):
                    if chart_kategori[ktg]=='~Industri':
                        chart_kategori[ktg]='Industri'
                chart_nilai = list(temp_df_[temp_df_.index.get_level_values('Aspek').isin([l_asp[asp]])]['~Average'].values)
                chart_nilai = chart_nilai[::-1]
                chart_nilai = [round(nl, 1) for nl in chart_nilai]

            elif (kriteria=='email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')) and (asp!=len(aspek_feas)-1):
                l_asp = temp_df_.index.get_level_values('Aspek').unique().to_list()
                chart_kategori = list(temp_df_[temp_df_.index.get_level_values('Aspek').isin([l_asp[asp]])].index.get_level_values('Brand').values)
                chart_kategori = chart_kategori[::-1]
                for ktg in range(len(chart_kategori)):
                    if chart_kategori[ktg]=='~Industri':
                        chart_kategori[ktg]='Industri'
                chart_nilai = list(temp_df_[temp_df_.index.get_level_values('Aspek').isin([l_asp[asp]])]['~Average'].values)
                chart_nilai = chart_nilai[::-1]
                chart_nilai = [round(nl, 1) for nl in chart_nilai]

            else:
                chart_kategori = list(temp_sub.index.get_level_values('Brand').values)
                chart_kategori = chart_kategori[::-1]
                for ktg in range(len(chart_kategori)):
                    if chart_kategori[ktg]=='~Industri':
                        chart_kategori[ktg]='Industri'
                chart_nilai = list(temp_sub['~Average'].values)
                chart_nilai = chart_nilai[::-1]
                chart_nilai = [round(nl, 1) for nl in chart_nilai]
            chart_nilai = ['#N/A' if str(x)=='nan' else x for x in chart_nilai]


            chart_data.categories = chart_kategori

            chart_data.add_series('Average', chart_nilai, '#,0"%"')

            # add chart to slide --------------------
            left, top, width, height = Inches(0.5), Inches(1.15), Inches(9), Inches(3.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
            ).chart

            chart.value_axis.maximum_scale = 100
            chart.value_axis.minimum_scale = 0
            chart.value_axis.major_unit = 20
            chart.value_axis.has_major_gridlines = False
            chart.value_axis.tick_labels.font.size = Pt(10)
            if kriteria=='call center':
                chart.category_axis.tick_labels.font.size = Pt(10)
            elif kriteria=='twitter' or kriteria=='online chat':
                chart.category_axis.tick_labels.font.size = Pt(6)
            elif kriteria=='email':
                chart.category_axis.tick_labels.font.size = Pt(4)
        #     chart.gap_width = 30

            chart.plots[0].has_data_label = True
            chart.plots[0].gap_width = 30
            for lbl in range(len(chart_nilai)):
                chart.plots[0].series[0].points[lbl].data_label.text_frame.text = str(chart_nilai[lbl])+"%"
                if kriteria=='call center':
                    chart.plots[0].series[0].points[lbl].data_label.font.size = Pt(10)
                elif kriteria=='twitter' or kriteria=='online chat':
                    chart.plots[0].series[0].points[lbl].data_label.font.size = Pt(6)
                elif kriteria=='email':
                    chart.plots[0].series[0].points[lbl].data_label.font.size = Pt(4)
                chart.plots[0].series[0].points[lbl].data_label.position = XL_LABEL_POSITION.INSIDE_END
                if chart_nilai[lbl]==0 or chart_nilai[lbl]=='#N/A':
                    chart.plots[0].series[0].points[lbl].data_label.font.color.rgb = RGBColor(255,255,255)
                chart.plots[0].series[0].points[lbl].format.fill.solid()
                if chart_kategori[lbl]==client:
                    chart.plots[0].series[0].points[lbl].format.fill.fore_color.rgb = RGBColor(0,112,192)
                    chart.plots[0].series[0].points[lbl].data_label.font.color.rgb = RGBColor(255,255,255)
                elif chart_kategori[lbl]=='Industri':
                    chart.plots[0].series[0].points[lbl].format.fill.fore_color.rgb = RGBColor(192,0,0)
                    chart.plots[0].series[0].points[lbl].data_label.font.color.rgb = RGBColor(255,255,255)
                else:
                    chart.plots[0].series[0].points[lbl].format.fill.fore_color.rgb = RGBColor(191,191,191)

            left = Inches(0.5)
            top = Inches(0.25)
            width = Inches(9)
            height = Inches(3)
            txBox = shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = ("Grafik dan Tabel Perbandingan Kinerja "
                      +client+" & Brand Lainnya\nHasil Pemantauan "+nama_bulan[0]+
                      " - "+nama_bulan[len(list_bulan)-1]+" "+str(data_tahun)+"\nUntuk KPI dan Atribut "+
                      subkpi_)
            p.font.bold = True
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER

            top = Inches(4.6)
            height = Inches(1.5)
            cols = 4
            rows = len(l_aspek_sub)+3

            if kriteria=='call center':
                if len(list_aspek_sub) > 1 and len(list_aspek_sub)!=len(l_aspek_sub):
                    if len(list_aspek_sub[-1]) != len(l_aspek_sub):
                        rows = len(l_aspek_sub)+2
                    elif (list_aspek_sub[-1]==l_aspek_sub).all()==True:
                        rows = len(l_aspek_sub)+3
                    else:
                        rows = len(l_aspek_sub)+2
            else:
                if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility') and (asp==1):
                    rows = len([l_aspek_sub])+3
                else:
                    rows = len(l_aspek_sub)+3

                if len(list_aspek_sub) > 1 and len(list_aspek_sub)!=len(l_aspek_sub):
                    if len(list_aspek_sub[-1]) != len(l_aspek_sub):
                        rows = len(l_aspek_sub)+2
                    elif (len(list_aspek_sub[-1]) == len(l_aspek_sub)) and len(l_aspek_sub)==1:
                        if (list_aspek_sub[-1]==l_aspek_sub)==True:
                            rows = len(l_aspek_sub)+3
                    elif (len(list_aspek_sub[-1]) != len(l_aspek_sub)) and len(l_aspek_sub)>1:
                        if (list_aspek_sub[-1]==l_aspek_sub).all()==True:
                            rows = len(l_aspek_sub)+3
                    else:
                        rows = len(l_aspek_sub)+2

            left = Inches(0.25)
            width = Inches(9.25)

            tabel_dibar = shapes.add_table(rows, cols, left, top, width, height).table

            tabel_dibar.columns[0].width = Inches(.25)
            tabel_dibar.columns[1].width = Inches(6.25)
            lebar_kolom_lainnya = Inches(1.5)

            for ke in range(cols-2):
                tabel_dibar.columns[ke+2].width = lebar_kolom_lainnya

            for baris in range(rows):
                list_temp = []
                if baris==0:
                    if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility'):
                        header = ['','ATRIBUT %s %s'%(str.upper(subkpi_), str(asp+1)),client, "Rata-rata Industri"]
                    elif (kriteria == 'email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_ == 'Soft Skill')) and (asp==0):
                        if subkpi_=='Service Standard Consistency':
                            header = ['','ATRIBUT %s - %s'%(str.upper(subkpi_), 'BAGIAN PEMBUKA'),client, "Rata-rata Industri"]
                        elif subkpi_=='Soft Skill':
                            header = ['','ATRIBUT %s - %s'%(str.upper(subkpi_), 'ETIKA YANG BAIK'),client, "Rata-rata Industri"]
                    elif (kriteria == 'email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_ == 'Soft Skill')) and (asp==1):
                        if subkpi_ =='Service Standard Consistency':
                            header = ['','ATRIBUT %s - %s'%(str.upper(subkpi_), 'BAGIAN PENUTUP'),client, "Rata-rata Industri"]
                        elif subkpi_=='Soft Skill':
                            header = ['','ATRIBUT %s - %s'%(str.upper(subkpi_), 'BAHASA SERVICE YANG BAIK'),client, "Rata-rata Industri"]
                    else:
                        header = ['','ATRIBUT '+str.upper(subkpi_),client, "Rata-rata Industri"]
                    for kolom in range(cols):
                        _set_cell_border(tabel_dibar.cell(baris,kolom))
                        tabel_dibar.cell(baris,kolom).text = header[kolom]
                        tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                        tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].font.size = Pt(11)
                        tabel_dibar.cell(baris,kolom).vertical_anchor = MSO_ANCHOR.MIDDLE
                        tabel_dibar.cell(baris,kolom).fill.solid()
                        tabel_dibar.cell(baris,kolom).fill.fore_color.rgb = RGBColor(155,187,89)
                        tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    tabel_dibar.cell(baris,0).merge(tabel_dibar.cell(baris,1))

                elif baris<len(l_aspek_sub)+1:
                    if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility') and (asp==0):
                        apa_asp = apa[apa.index.get_level_values('Aspek').isin([l_aspek_sub[baris-1]])]
                        value_avg_client = apa_asp[apa_asp.index.get_level_values('Brand').isin([client])].values[0][-1]
                        value_avg_industri = apa_asp[apa_asp.index.get_level_values('Brand').isin(['~Industri'])].values[0][-1]
                    elif (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility') and (asp==1):
                        temp_df_ = data_aspek[data_aspek.index.get_level_values('SUB KPI').isin([subkpi_])]
                        value_client = temp_df_[temp_df_.index.get_level_values('Aspek').isin([list_aspek[0]])]
                        value_avg_client = value_client[value_client.index.get_level_values('Brand').isin([client])].values[0][-1]
                        value_avg_industri = value_client[value_client.index.get_level_values('Brand').isin(['~Industri'])].values[0][-1]

                    elif (kriteria=='email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')) and (asp!=len(aspek_feas)-1):
                        if asp==0:
                            apa_asp = apa1[apa1.index.get_level_values('Aspek').isin([l_aspek_sub[baris-1]])]
                        else:
                            apa_asp = apa2[apa2.index.get_level_values('Aspek').isin([l_aspek_sub[baris-1]])]
                        value_avg_client = apa_asp[apa_asp.index.get_level_values('Brand').isin([client])].values[0][-1]
                        value_avg_industri = apa_asp[apa_asp.index.get_level_values('Brand').isin(['~Industri'])].values[0][-1]
                    else:
                        temp_df_asp = temp_df_[temp_df_.index.get_level_values('Aspek').isin([l_aspek_sub[baris-1]])]
                        value_avg_client = temp_df_asp[temp_df_asp.index.get_level_values('Brand').isin([client])].values[0][-1]
                        value_avg_industri = temp_df_asp[temp_df_asp.index.get_level_values('Brand').isin(['~Industri'])].values[0][-1]

                    if len(list_aspek_sub) > 1:
                        pjg = len(list_aspek_sub)
                        for sub_tabel in range(pjg):
                            if l_aspek_sub[baris-1] in list_aspek_sub[sub_tabel]:
                                list_temp.extend([baris+(len(list_aspek_sub[sub_tabel-1])*sub_tabel)])
                    else:
                        list_temp.extend([baris])

                    list_temp.extend([l_aspek_sub[baris-1]])
                    list_temp.extend([value_avg_client])
                    list_temp.extend([value_avg_industri])

                    for kolom in range(cols):
                        t_value = list_temp[kolom]
                        if str(t_value)=='nan':
                            val_ = '#N/A'
                        elif type(t_value)==np.float64:
                            val_ = str(round(t_value,1))+"%"
                        else:
                            val_ = str(t_value)
                        _set_cell_border(tabel_dibar.cell(baris,kolom))
                        tabel_dibar.cell(baris,kolom).text = val_
                        tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                        tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].font.size = Pt(10)

                        if kolom==1:
                            tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                        else:
                            tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            tabel_dibar.cell(baris,kolom).vertical_anchor = MSO_ANCHOR.MIDDLE

                        tabel_dibar.cell(baris,kolom).fill.solid()
                        if baris%2!=0:
                            tabel_dibar.cell(baris,kolom).fill.fore_color.rgb = RGBColor(255, 128, 0)
                        else:
                            tabel_dibar.cell(baris,kolom).fill.fore_color.rgb = RGBColor(255, 255, 153)

                elif baris==len(l_aspek_sub)+1:

                    if (rows-2!=baris) and (laspek!=len(list_aspek_sub)-1):
                        for kolom in range(cols):
                            _set_cell_border(tabel_dibar.cell(baris,kolom))
                            if kolom==0:
                                tabel_dibar.cell(baris,kolom).text = str(baris+(4*laspek))
                                tabel_dibar.cell(baris,kolom).vertical_anchor = MSO_ANCHOR.MIDDLE
                                tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            elif kolom == 1:
                                tabel_dibar.cell(baris,kolom).text = '...'
                                tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                            else:
                                tabel_dibar.cell(baris,kolom).text = ''
                                tabel_dibar.cell(baris,kolom).vertical_anchor = MSO_ANCHOR.MIDDLE
                                tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                            tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].font.size = Pt(10)
                            tabel_dibar.cell(baris,kolom).fill.solid()
                            tabel_dibar.cell(baris,kolom).fill.fore_color.rgb = RGBColor(255, 255, 255)

                    else:
                        if (kriteria == 'twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility'):
                            temp_df_1 = data_aspek[data_aspek.index.get_level_values('SUB KPI').isin([subkpi_])]
                            if asp==0:
                                value_sub_client = temp_df_1[temp_df_1.index.get_level_values('Brand').isin([client])].values[0][-1]
                                value_sub_industri = temp_df_1[temp_df_1.index.get_level_values('Brand').isin(['~Industri'])].values[0][-1]
                            else:
                                value_sub_client = temp_df_1[temp_df_1.index.get_level_values('Brand').isin([client])].values[1][-1]
                                value_sub_industri = temp_df_1[temp_df_1.index.get_level_values('Brand').isin(['~Industri'])].values[1][-1]
                        elif (kriteria=='email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')) and (asp!=len(aspek_feas)-1):
                            temp_df_1 = data_aspek[data_aspek.index.get_level_values('SUB KPI').isin([subkpi_])]
                            if asp==0:
                                if subkpi_ == 'Service Standard Consistency':
                                    temp_df_1 = temp_df_1[temp_df_1.index.get_level_values('Aspek').str.contains('pembuka')]
                                elif subkpi_ == 'Soft Skill':
                                    temp_df_1 = temp_df_1[temp_df_1.index.get_level_values('Aspek').str.contains('etika yang baik')]
                            else:
                                if subkpi_ == 'Service Standard Consistency':
                                    temp_df_1 = temp_df_1[temp_df_1.index.get_level_values('Aspek').str.contains('penutup')]
                                elif subkpi_ == 'Soft Skill':
                                    temp_df_1 = temp_df_1[temp_df_1.index.get_level_values('Aspek').str.contains('Service yang baik')]

                            value_sub_client = temp_df_1[temp_df_1.index.get_level_values('Brand').isin([client])].values[0][-1]
                            value_sub_industri = temp_df_1[temp_df_1.index.get_level_values('Brand').isin(['~Industri'])].values[0][-1]
                        else:
                            value_sub_client = temp_sub[temp_sub.index.get_level_values('Brand').isin([client])].values[0][-1]
                            value_sub_industri = temp_sub[temp_sub.index.get_level_values('Brand').isin(['~Industri'])].values[0][-1]

                        list_temp.extend([''])
                        list_temp.extend(['TOTAL SCORE'])
                        list_temp.extend([value_sub_client])
                        list_temp.extend([value_sub_industri])

                        for kolom in range(cols):
                            t_value = list_temp[kolom]
                            if str(t_value)=='nan':
                                val_ = '#N/A'
                            elif type(t_value)==np.float64:
                                val_ = str(round(t_value,1))+"%"
                            else:
                                val_ = str(t_value)
                            _set_cell_border(tabel_dibar.cell(baris,kolom))
                            tabel_dibar.cell(baris,kolom).text = val_
                            tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
                            tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].font.size = Pt(10)
                            tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].font.bold = True
                            tabel_dibar.cell(baris,kolom).vertical_anchor = MSO_ANCHOR.MIDDLE
                            tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            tabel_dibar.cell(baris,kolom).fill.solid()
                            tabel_dibar.cell(baris,kolom).fill.fore_color.rgb = RGBColor(255, 0, 0)

                            if kolom >=2 :
                                tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            else:
                                tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            #             tabel_dibar[i].cell(j,0).merge(list_kinerja_akses[i].cell(j,1))

                else:
                    list_temp.extend([''])
                    list_temp.extend(['n Sample'])
                    for kolom in list_bulan:
                        list_temp.extend([data_mentah.Periode.nunique()])
                    for kolom in range(cols):
                        t_value = list_temp[kolom]
                        if type(t_value)==np.float64:
                            val_ = str(t_value)+"%"
                        elif kolom==cols-2:
                            val_ = str(t_value*len(list_bulan))
                        elif kolom==cols-1:
                            if (kriteria=='call center') and (subkpi_ in subkpi_ivr):
                                brand_ivr, brand_nonivr = od.get_ivr_nonivr(data_mentah, list_brand)
                                jumlah_ivr, jumlah_nonivr = len(brand_ivr), len(brand_nonivr)
                                if (client in brand_ivr):
                                    val_ = str(t_value*len(list_bulan)*jumlah_ivr)
                            else:
                                val_ = str(t_value*len(list_bulan)*len(list_brand))
                        else:
                            val_ = str(t_value)
                        _set_cell_border(tabel_dibar.cell(baris,kolom))
                        tabel_dibar.cell(baris,kolom).text = val_
                        tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                        tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].font.size = Pt(10)
                        tabel_dibar.cell(baris,kolom).vertical_anchor = MSO_ANCHOR.MIDDLE
                        tabel_dibar.cell(baris,kolom).fill.solid()
                        tabel_dibar.cell(baris,kolom).fill.fore_color.rgb = RGBColor(255, 255, 255)
                        if kolom >=2 :
                            tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        else:
                            tabel_dibar.cell(baris,kolom).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        #             tabel_dibar[i].cell(j,0).merge(list_kinerja_akses[i].cell(j,1))

            footer(gambar_footer, slide)


def plot_tabel_engagement(template_ppt, gambar_footer, list_bulan, data_xls, client, data_tahun):
    """
    Membuat slide tabel indeks engagement

    Parameter
    ---------
    template_ppt : template ptt yang digunakan
    gambar_footer : nama file gambar yang akan dipakai
    list_bulan : list array bulan yang ada di dalam dataframe
    data_xls : file data yang akan ditampilkan di slide
    client : nama client yang dibuatkan laporan
    data_tahun : tahun pembuatan laporan

    Return
    ------
    None
    """

    slide = template_ppt.slide_layouts[6]
    slide = template_ppt.slides.add_slide(slide)
    shapes = slide.shapes

    #judul
    left = Inches(0.5)
    top = Inches(0.25)
    width = Inches(9)
    height = Inches(3)
    txBox = shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "ENGAGEMENT INDEX\nHasil Pemantauan "+list_bulan[0]+" - "+list_bulan[-1]+" "+str(data_tahun)
    p.font.bold = True
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER

    tabel_be_avg = pd.read_excel(data_xls, sheet_name='engagement index aspek',index_col=[0,1,2,3,4,5,6], engine='openpyxl')
    list_aspek = tabel_be_avg.index.droplevel([0,1,2,3,4,6]).unique().values
    eng_index = tabel_be_avg.copy()
    eng_index.index = eng_index.index.droplevel([0,1,2,3,4])

    df_eng_final_pivot = pd.read_excel(data_xls, sheet_name = 'engagement index final', index_col=[0], engine='openpyxl')

    top = Inches(1.25)
    cols = 3
    rows = len(list_aspek)+2
    left = Inches(0.25)
    width = Inches(9.25)
    height = Inches(5)

    tb_engagement = shapes.add_table(rows, cols, left, top, width, height).table
    tb_engagement.columns[0].width = Inches(.5)
    tb_engagement.columns[1].width = Inches(7.5)
    tb_engagement.columns[2].width = Inches(1.5)

    for j in range(rows):
        list_temp = []
        if j==0:
            header = ['','ATRIBUT ENGAGEMENT INDEX',client]
            list_temp = header

            for k in range(cols):
                _set_cell_border(tb_engagement.cell(j,k))
                tb_engagement.cell(j,k).text = list_temp[k]
                tb_engagement.cell(j,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                tb_engagement.cell(j,k).text_frame.paragraphs[0].font.size = Pt(11)
                tb_engagement.cell(j,k).vertical_anchor = MSO_ANCHOR.MIDDLE
                tb_engagement.cell(j,k).fill.solid()
                tb_engagement.cell(j,k).fill.fore_color.rgb = RGBColor(155,187,89)
                tb_engagement.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tb_engagement.cell(j,0).merge(tb_engagement.cell(j,1))

        elif j<len(list_aspek)+1:
            l_aspek = list_aspek[j-1]
            value_aspek = eng_index.loc[(list_aspek[j-1],client),'~Average']

            list_temp.extend([j])
            list_temp.extend([list_aspek[j-1]])
            list_temp.extend([value_aspek])

            for k in range(cols):
                t_value = list_temp[k]
                if str(t_value)=='nan':
                    val_ = '#N/A'
                elif type(t_value)==np.float64:
                    t_value = round(t_value,1)
                    val_ = str(t_value)+"%"
                else:
                    val_ = str(t_value)
                _set_cell_border(tb_engagement.cell(j,k))
                tb_engagement.cell(j,k).text = val_
                tb_engagement.cell(j,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                tb_engagement.cell(j,k).text_frame.paragraphs[0].font.size = Pt(10)

                if k==1:
                    tb_engagement.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                else:
                    tb_engagement.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    tb_engagement.cell(j,k).vertical_anchor = MSO_ANCHOR.MIDDLE

                tb_engagement.cell(j,k).fill.solid()
                if j%2!=0:
                    tb_engagement.cell(j,k).fill.fore_color.rgb = RGBColor(255, 128, 0)
                else:
                    tb_engagement.cell(j,k).fill.fore_color.rgb = RGBColor(255, 255, 153)

        else:
            list_temp.extend([''])
            list_temp.extend(['TOTAL SCORE'])
            list_temp.extend([df_eng_final_pivot.loc[client,'~Average']])

            for k in range(cols):
                t_value = list_temp[k]
                if str(t_value)=='nan':
                    val_ = '#N/A'
                elif type(t_value)==np.float64:
                    t_value = round(t_value,1)
                    val_ = str(t_value)+"%"
                else:
                    val_ = str(t_value)
                _set_cell_border(tb_engagement.cell(j,k))
                tb_engagement.cell(j,k).text = val_
                tb_engagement.cell(j,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
                tb_engagement.cell(j,k).text_frame.paragraphs[0].font.size = Pt(10)
                tb_engagement.cell(j,k).text_frame.paragraphs[0].font.bold = True
                tb_engagement.cell(j,k).vertical_anchor = MSO_ANCHOR.MIDDLE
                tb_engagement.cell(j,k).fill.solid()
                tb_engagement.cell(j,k).fill.fore_color.rgb = RGBColor(255, 0, 0)
                if k >=2 :
                    tb_engagement.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                else:
                    tb_engagement.cell(j,k).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT



def plot_tabel_perbaikan(template_ppt, gambar_footer, list_bulan, data_xls, bobot, client, data_tahun, kriteria):
    """
    Membuat slide tabel perbaikan

    Parameter
    ---------
    template_ppt : template ptt yang digunakan
    gambar_footer : nama file gambar yang akan dipakai
    list_bulan : list array bulan yang ada di dalam dataframe
    data_xls : file data yang akan ditampilkan di slide
    bobot : dataframe yang berisi daftar bobot perhitungan indeks
    client : nama client yang dibuatkan laporan
    data_tahun : tahun pembuatan laporan
    kriteria : kategori SEI yang akan dibuat ('call center', 'online chat', 'twitter', 'email')

    Return
    ------
    None
    """

    data_subkpi = pd.read_excel(data_xls, sheet_name='subkpi', engine='openpyxl')
    data_subkpi['Dimensi'] = data_subkpi['Dimensi'].fillna(method='ffill')
    data_subkpi['KPI'] = data_subkpi['KPI'].fillna(method='ffill')
    data_subkpi['SUB KPI'] = data_subkpi['SUB KPI'].fillna(method='ffill')
    data_subkpi = data_subkpi.set_index(['Dimensi','KPI','SUB KPI','Brand'])
    list_subkpi = data_subkpi.index.get_level_values('SUB KPI').unique()

    df_aspek_perbaikan = pd.read_excel(data_xls, sheet_name = 'aspek perbaikan', index_col=[0,1,2,3,4,5], engine='openpyxl')
    bobot_indeks_ = pd.read_excel(bobot, sheet_name='perbaikan', engine='openpyxl')
    bobot_indeks_ = bobot_indeks(bobot_indeks_, perbaikan=True)

    data_mentah = pd.read_excel(data_xls, sheet_name='data mentah', engine='openpyxl')
    list_brand = data_mentah.Brand.unique()
    if (kriteria=='call center'):
        brand_ivr, brand_nonivr = od.get_ivr_nonivr(data_mentah, list_brand)


    for j_subkpi in range(len(list_subkpi)):
        subkpi_ = list_subkpi[j_subkpi]
        if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility'):
            asp = 2
        elif (kriteria=='email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')):
            asp = 3
        else:
            asp = 1

        for loop in range(asp):
            temp_df_ = df_aspek_perbaikan.unstack().stack(level=0, dropna=False).unstack()
            temp_df_ = temp_df_[(client)].reset_index()
            if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility'):
                temp_df_ = temp_df_[(temp_df_['SUB KPI'].isin([subkpi_])) & (temp_df_['NOMOR LEVEL 1'].isin([loop+1]))]
                var = bobot_indeks_[['ASPEK LAPORAN']][:4]
            elif (kriteria=='email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')):
                if loop==0:
                    if subkpi_=='Service Standard Consistency':
                        temp_df_ = temp_df_[(temp_df_['SUB KPI'].isin([subkpi_])) & (temp_df_['ASPEK LAPORAN'].str.contains('Pembuka'))]
                    elif subkpi_ == 'Soft Skill':
                        temp_df_ = temp_df_[(temp_df_['SUB KPI'].isin([subkpi_])) & (temp_df_['ASPEK LAPORAN'].str.contains('Etika yang Baik'))]
                elif loop==1:
                    if subkpi_=='Service Standard Consistency':
                        temp_df_ = temp_df_[(temp_df_['SUB KPI'].isin([subkpi_])) & (temp_df_['ASPEK LAPORAN'].str.contains('Penutup'))]
                    elif subkpi_ == 'Soft Skill':
                        temp_df_ = temp_df_[(temp_df_['SUB KPI'].isin([subkpi_])) & (temp_df_['ASPEK LAPORAN'].str.contains('Bahasa Service yang Baik'))]
                else:
                    if subkpi_=='Service Standard Consistency':
                        temp_df_ = temp_df_[(temp_df_['SUB KPI'].isin([subkpi_])) & ~(temp_df_['ASPEK LAPORAN'].str.contains('Pembuka'))]
                        temp_df_ = temp_df_[~(temp_df_['ASPEK LAPORAN'].str.contains('Penutup'))]
                    elif subkpi_ == 'Soft Skill':
                        temp_df_ = temp_df_[(temp_df_['SUB KPI'].isin([subkpi_])) & ~(temp_df_['ASPEK LAPORAN'].str.contains('Etika yang Baik'))]
                        temp_df_ = temp_df_[~(temp_df_['ASPEK LAPORAN'].str.contains('Bahasa Service yang Baik'))]
                asp_ = temp_df_['ASPEK LAPORAN'].unique()#.to_list()
                if loop==0:
                    subaspek_1 = asp_
                elif loop==1:
                    subaspek_2 = asp_
            else:
                temp_df_ = temp_df_[temp_df_['SUB KPI'].isin([subkpi_])]
            list_aspek = temp_df_['ASPEK LAPORAN'].values


            slide_perbaikan = template_ppt.slide_layouts[6]
            slide = template_ppt.slides.add_slide(slide_perbaikan)
            shapes = slide.shapes

            #judul
            left = Inches(0.5)
            top = Inches(0.25)
            width = Inches(9)
            height = Inches(3)
            txBox = shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = ("AREA PERBAIKAN\nHasil Pemantauan "+
                      list_bulan[0]+" - "+list_bulan[-1]+" "+str(data_tahun))
            p.font.bold = True
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER

            top = Inches(1.25)
            height = Inches(1.5)
            cols = 5
            rows = len(list_aspek)+1
            left = Inches(0.25)
            width = Inches(9)

            list_perbaikan = []
            list_perbaikan = shapes.add_table(rows, cols, left, top, width, height).table
            list_perbaikan.columns[0].width = Inches(.4)
            list_perbaikan.columns[1].width = Inches(3.75)
            list_perbaikan.rows[0].height = Inches(0.5)

            for baris in range(rows):
                list_temp = []
                if baris==0:
                    if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility'):
                        header = ['','ATRIBUT %s %s'%(str.upper(subkpi_), str(loop+1)),'Frekuensi Perbaikan','Impact Index','Priority Index']
                    elif (kriteria=='email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')):
                        if loop==0:
                            if subkpi_=='Service Standard Consistency':
                                header = ['','ATRIBUT %s - %s'%(str.upper(subkpi_), 'BAGIAN PEMBUKA'),'Frekuensi Perbaikan','Impact Index','Priority Index']
                            elif subkpi_=='Soft Skill':
                                header = ['','ATRIBUT %s - %s'%(str.upper(subkpi_), 'ETIKA YANG BAIK'),'Frekuensi Perbaikan','Impact Index','Priority Index']
                        elif loop==1:
                            if subkpi_=='Service Standard Consistency':
                                header = ['','ATRIBUT %s - %s'%(str.upper(subkpi_), 'BAGIAN PENUTUP'),'Frekuensi Perbaikan','Impact Index','Priority Index']
                            elif subkpi_=='Soft Skill':
                                header = ['','ATRIBUT %s - %s'%(str.upper(subkpi_), 'BAHASA SERVICE YANG BAIK'),'Frekuensi Perbaikan','Impact Index','Priority Index']
                        else:
                            header = ['','ATRIBUT '+str.upper(subkpi_),'Frekuensi Perbaikan','Impact Index','Priority Index']
                    else:
                        header = ['','ATRIBUT '+str.upper(subkpi_),'Frekuensi Perbaikan','Impact Index','Priority Index']
                    list_temp = header

                    for k in range(cols):
                        _set_cell_border(list_perbaikan.cell(baris,k))
                        list_perbaikan.cell(baris,k).text = list_temp[k]
                        list_perbaikan.cell(baris,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                        list_perbaikan.cell(baris,k).text_frame.paragraphs[0].font.size = Pt(11)
                        list_perbaikan.cell(baris,k).vertical_anchor = MSO_ANCHOR.MIDDLE
                        list_perbaikan.cell(baris,k).fill.solid()
                        list_perbaikan.cell(baris,k).fill.fore_color.rgb = RGBColor(155,187,89)
                        list_perbaikan.cell(baris,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    list_perbaikan.cell(baris,0).merge(list_perbaikan.cell(baris,1))

                else:
                    if (kriteria=='twitter' or kriteria=='online chat' or kriteria=='email') and (subkpi_ == 'Feasibility') and (loop==0):
                        aspek_ = var['ASPEK LAPORAN'][baris-1]
                    elif (kriteria=='email') and ((subkpi_ == 'Service Standard Consistency') or (subkpi_=='Soft Skill')):
                        if loop==0:
                            aspek_ = subaspek_1[baris-1]
                        elif loop==1:
                            aspek_ = subaspek_2[baris-1]
                        else:
                            aspek_ = asp_[baris-1]
                    else:
                        aspek_ = list_aspek[baris-1]
                    frek_ = temp_df_['Frekuensi Perbaikan'].reset_index(drop=True)[baris-1]
                    prindex_ = temp_df_['_Priority Index'].reset_index(drop=True)[baris-1]

                    bobot_perbaikan = pd.read_excel(bobot, sheet_name = 'perbaikan', engine='openpyxl')
                    if kriteria=='call center':
                        if client in brand_ivr:
                            imindex_ = bobot_perbaikan[bobot_perbaikan['ASPEK LAPORAN'].isin([list_aspek[baris-1]])]['BOBOT IMPACT IVR'].values[0]
                        else:
                            imindex_ = bobot_perbaikan[bobot_perbaikan['ASPEK LAPORAN'].isin([list_aspek[baris-1]])]['BOBOT IMPACT NON-IVR'].values[0]
                    else:
                        imindex_ = bobot_perbaikan[bobot_perbaikan['ASPEK LAPORAN'].isin([list_aspek[baris-1]])]['BOBOT IMPACT'].values[0]

                    list_temp.extend([baris])
                    list_temp.extend([aspek_])
                    list_temp.extend([frek_])
                    list_temp.extend([round(imindex_,3)])
                    list_temp.extend([round(prindex_,3)])

                    #jika data yang ada maksimal hanya terisi sampai bulan Juni, pakai data ini
                    for k in range(cols):
                        t_value = list_temp[k]
                        if k > 1 and len(str(t_value))<5:
                            if k==2:
                                val_ = str(int(t_value))
                            else:
                                val_ = str(t_value)+("0"*(5-len(str(t_value))))
                        else:
                            val_ = str(t_value)
                        _set_cell_border(list_perbaikan.cell(baris,k))
                        list_perbaikan.cell(baris,k).text = val_
                        list_perbaikan.cell(baris,k).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                        list_perbaikan.cell(baris,k).text_frame.paragraphs[0].font.size = Pt(10)

                        if k==1:
                            list_perbaikan.cell(baris,k).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                        else:
                            list_perbaikan.cell(baris,k).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            list_perbaikan.cell(baris,k).vertical_anchor = MSO_ANCHOR.MIDDLE

                        list_perbaikan.cell(baris,k).fill.solid()
                        if baris%2!=0:
                            list_perbaikan.cell(baris,k).fill.fore_color.rgb = RGBColor(255, 128, 0)
                        else:
                            list_perbaikan.cell(baris,k).fill.fore_color.rgb = RGBColor(255, 255, 153)

            footer(gambar_footer, slide)
